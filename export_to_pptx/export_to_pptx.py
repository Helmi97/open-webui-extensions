"""
title: Export to PPTX
description: Export Assistant Message to PPTX
version: 1.0.0
author: Helmi Chaouachi
git_url: https://github.com/Helmi97/open-webui-extensions/tree/main/export_to_pptx
icon_url: https://www.svgrepo.com/show/373991/powerpoint2.svg
required_open_webui_version: 0.8.0
requirements: python-pptx,markdown,requests,beautifulsoup4,pillow
"""

from __future__ import annotations

import base64
import html
import io
import json
import logging
import math
import re
import time
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Optional
from urllib.parse import urlparse

import markdown
import requests
from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from open_webui.models.users import UserModel
from open_webui.utils.chat import generate_chat_completion

LOGGER = logging.getLogger(__name__)

MERMAID_BLOCK_RE = re.compile(
    r"```mermaid[ \t]*\n(.*?)\n```",
    flags=re.IGNORECASE | re.DOTALL,
)
HEADING_LINE_RE = re.compile(r"^\s*#\s+(.+?)\s*$", flags=re.MULTILINE)
MARKDOWN_IMAGE_RE = re.compile(r"!\[(?P<alt>[^\]]*)\]\((?P<src>[^)]+)\)")
HTML_IMG_RE = re.compile(r"(?is)<img\b[^>]*>")
PLACEHOLDER_RE = re.compile(r"\{\{\s*([A-Za-z0-9_]+)\s*\}\}")
BAR_LOGO_PLACEHOLDER_RE = re.compile(
    r"(\{\{\s*(HEADER_LOGO|FOOTER_LOGO)\s*\}\})",
    flags=re.IGNORECASE,
)
WRAPPED_CODE_FENCE_RE = re.compile(
    r"^\s*```(?:markdown|md)?\s*\n(?P<body>.*)\n```\s*$",
    flags=re.IGNORECASE | re.DOTALL,
)
SENTENCE_SPLIT_RE = re.compile(r"(?<=[.!?])\s+")

DEFAULT_LLM_PROMPT = (
    "You convert an assistant message into a clean markdown document for PowerPoint "
    "export. Stay as faithful as possible to the original content. Write a document, "
    "not presentation bullets. Use markdown headings with exactly one H1 title, then "
    "H2 section titles, then H3/H4 subsections when useful. Keep paragraphs readable. "
    "Use bullet lists or numbered lists only when they improve clarity. Preserve any "
    "special visual tokens exactly as provided, without changing, deleting, or merging "
    "them. Place each visual token on its own line near the most relevant text. Return "
    "only markdown."
)


@dataclass(slots=True)
class AssetToken:
    token: str
    src: str
    alt: str
    width_px: int = 0
    height_px: int = 0


@dataclass(slots=True)
class OutlineNode:
    level: int
    title: str
    blocks: list[str] = field(default_factory=list)
    children: list["OutlineNode"] = field(default_factory=list)


@dataclass(slots=True)
class ContentBlock:
    kind: str
    text: str = ""
    items: list[str] = field(default_factory=list)
    ordered: bool = False
    src: str = ""
    alt: str = ""
    rows: list[list[str]] = field(default_factory=list)


@dataclass(slots=True)
class SlideSpec:
    kind: str
    title: str = ""
    section_title: str = ""
    subsection_title: str = ""
    blocks: list[ContentBlock] = field(default_factory=list)
    image: Optional[ContentBlock] = None


class Action:
    class Valves(BaseModel):
        debug: bool = Field(
            default=False,
            description="Enable verbose server-side debug logging for this action.",
        )
        priority: int = Field(
            default=0,
            description="Controls button display order (lower = appears first).",
        )
        filename_prefix: str = Field(
            default="message",
            description="Prefix used for the downloaded PPTX file name.",
        )
        use_llm: bool = Field(
            default=True,
            description=(
                "Use the current chat model to restructure the message into a clean "
                "markdown document before rendering slides."
            ),
        )
        llm_model: str = Field(
            default="",
            description=(
                "Optional model ID used for document restructuring. Defaults to the "
                "current chat model."
            ),
        )
        llm_temperature: float = Field(
            default=0.2,
            ge=0.0,
            le=2.0,
            description="Temperature used when restructuring the document.",
        )
        llm_max_tokens: int = Field(
            default=4096,
            ge=256,
            le=16384,
            description="Maximum tokens used for the document restructuring request.",
        )
        llm_prompt: str = Field(
            default=DEFAULT_LLM_PROMPT,
            description="System prompt used when `use_llm` is enabled.",
        )
        request_timeout_s: int = Field(
            default=30,
            description="HTTP timeout in seconds for downloading images.",
        )
        mermaid_scale: int = Field(
            default=2,
            description=(
                "Rasterization scale used when converting rendered Mermaid diagrams "
                "from the browser into PNG images."
            ),
        )
        slide_width_in: float = Field(
            default=13.333,
            description="Presentation width in inches. 13.333 is 16:9 widescreen.",
        )
        slide_height_in: float = Field(
            default=7.5,
            description="Presentation height in inches. 7.5 is 16:9 widescreen.",
        )
        content_margin_in: float = Field(
            default=0.45,
            description="Base inner margin used for slide content.",
        )
        header_height_in: float = Field(
            default=0.35,
            description="Header bar height in inches. Use 0 to disable the header.",
        )
        footer_height_in: float = Field(
            default=0.35,
            description="Footer bar height in inches. Use 0 to disable the footer.",
        )
        header_text_left: str = Field(
            default="{{ PRESENTATION_TITLE }}",
            description="Header left text. Supports placeholders.",
        )
        header_text_middle: str = Field(
            default="",
            description="Header middle text. Supports placeholders.",
        )
        header_text_right: str = Field(
            default="{{ CURRENT_PAGE }}/{{ TOTAL_PAGES }}",
            description="Header right text. Supports placeholders.",
        )
        footer_text_left: str = Field(
            default="{{ FILE_NAME }}",
            description="Footer left text. Supports placeholders.",
        )
        footer_text_middle: str = Field(
            default="{{ EXPORT_DATE }} {{ EXPORT_TIME }}",
            description="Footer middle text. Supports placeholders.",
        )
        footer_text_right: str = Field(
            default="{{ USER_NAME }}",
            description="Footer right text. Supports placeholders.",
        )
        header_logo_url: str = Field(
            default="",
            description=(
                "Logo URL or data URL used when `{{ HEADER_LOGO }}` appears in a "
                "header field."
            ),
        )
        footer_logo_url: str = Field(
            default="",
            description=(
                "Logo URL or data URL used when `{{ FOOTER_LOGO }}` appears in a "
                "footer field."
            ),
        )
        header_bg_color: str = Field(
            default="#F1F5F9",
            description="Header background color.",
        )
        footer_bg_color: str = Field(
            default="#F8FAFC",
            description="Footer background color.",
        )
        header_text_color: str = Field(
            default="#334155",
            description="Header text color.",
        )
        footer_text_color: str = Field(
            default="#475569",
            description="Footer text color.",
        )
        title_slide_subtitle: str = Field(
            default="Prepared {{ EXPORT_DATE }} {{ EXPORT_TIME }}",
            description="Subtitle shown on the title slide. Supports placeholders.",
        )
        title_slide_bg_color: str = Field(
            default="#0F172A",
            description="Title slide background color.",
        )
        title_slide_bg_image_url: str = Field(
            default="",
            description=(
                "Optional background image URL or data URL for the title slide."
            ),
        )
        section_title_slide_bg_color: str = Field(
            default="#E2E8F0",
            description="Section title slide background color.",
        )
        section_title_slide_bg_image_url: str = Field(
            default="",
            description=(
                "Optional background image URL or data URL for section title slides."
            ),
        )
        content_slide_bg_color: str = Field(
            default="#FFFFFF",
            description="Normal content slide background color.",
        )
        title_font_name: str = Field(
            default="Aptos Display",
            description="Main title font.",
        )
        title_font_size_pt: int = Field(
            default=28,
            description="Main title font size.",
        )
        title_color: str = Field(
            default="#F8FAFC",
            description="Main title color.",
        )
        subtitle_font_name: str = Field(
            default="Aptos",
            description="Subtitle font.",
        )
        subtitle_font_size_pt: int = Field(
            default=16,
            description="Subtitle font size.",
        )
        subtitle_color: str = Field(
            default="#CBD5E1",
            description="Subtitle color.",
        )
        section_title_font_name: str = Field(
            default="Aptos Display",
            description="Section title slide font.",
        )
        section_title_font_size_pt: int = Field(
            default=26,
            description="Section title slide font size.",
        )
        section_title_color: str = Field(
            default="#0F172A",
            description="Section title slide text color.",
        )
        body_title_font_name: str = Field(
            default="Aptos Display",
            description="Content slide title font.",
        )
        body_title_font_size_pt: int = Field(
            default=22,
            description="Content slide title font size.",
        )
        body_title_color: str = Field(
            default="#0F172A",
            description="Content slide title color.",
        )
        body_font_name: str = Field(
            default="Aptos",
            description="Body font.",
        )
        body_font_size_pt: int = Field(
            default=16,
            description="Body font size.",
        )
        body_text_color: str = Field(
            default="#1E293B",
            description="Body text color.",
        )
        accent_color: str = Field(
            default="#2563EB",
            description="Accent color used in grids and separators.",
        )
        bullet_grid_bg_color: str = Field(
            default="#DBEAFE",
            description="Background color used for the short bullet grid layout.",
        )
        bullet_grid_text_color: str = Field(
            default="#1E3A8A",
            description="Text color used for the short bullet grid layout.",
        )
        closing_slide_enabled: bool = Field(
            default=True,
            description="Append a final thank-you / goodbye slide at the end.",
        )
        closing_slide_title: str = Field(
            default="Thank You",
            description="Main text shown on the final slide. Supports placeholders.",
        )
        closing_slide_subtitle: str = Field(
            default="Goodbye and see you soon.",
            description="Subtitle shown on the final slide. Supports placeholders.",
        )
        closing_slide_bg_color: str = Field(
            default="#0F172A",
            description="Final slide background color.",
        )
        closing_slide_bg_image_url: str = Field(
            default="",
            description="Optional background image URL or data URL for the final slide.",
        )
        closing_slide_title_color: str = Field(
            default="#F8FAFC",
            description="Final slide title color.",
        )
        closing_slide_subtitle_color: str = Field(
            default="#CBD5E1",
            description="Final slide subtitle color.",
        )

    def __init__(self):
        self.valves = self.Valves()
        self._image_info_cache: dict[str, tuple[bytes, int, int]] = {}

    def _preview(self, value: Any, limit: int = 200) -> str:
        text = value if isinstance(value, str) else str(value)
        text = text.replace("\r", "\\r").replace("\n", "\\n")
        if len(text) > limit:
            return f"{text[:limit]}...<truncated {len(text) - limit} chars>"
        return text

    def _debug_log(self, message: str, **context: Any) -> None:
        if not self.valves.debug:
            return

        if not context:
            LOGGER.info("[export_to_pptx] %s", message)
            return

        rendered_parts: list[str] = []
        for key, value in context.items():
            if isinstance(value, (dict, list, tuple)):
                try:
                    rendered = json.dumps(value, default=str)
                except TypeError:
                    rendered = str(value)
            else:
                rendered = str(value)
            rendered_parts.append(f"{key}={self._preview(rendered)}")

        LOGGER.info("[export_to_pptx] %s | %s", message, " | ".join(rendered_parts))

    def _summarize_slide_specs(self, slide_specs: list[SlideSpec]) -> dict[str, int]:
        summary: dict[str, int] = {}
        for spec in slide_specs:
            summary[spec.kind] = summary.get(spec.kind, 0) + 1
        return summary

    def build_filename(self, message_id: str) -> str:
        safe_id = re.sub(r"[^A-Za-z0-9._-]+", "-", (message_id or "").strip())
        safe_id = safe_id.strip("-._") or "message"
        prefix = re.sub(r"[^A-Za-z0-9._-]+", "-", self.valves.filename_prefix).strip(
            "-._"
        )
        prefix = prefix or "message"
        return f"{prefix}-{safe_id}.pptx"

    async def emit_status(
        self,
        description: str,
        done: bool = False,
        __event_emitter__=None,
    ):
        if __event_emitter__:
            await __event_emitter__(
                {
                    "type": "status",
                    "data": {"description": description, "done": done},
                }
            )

    async def emit_error(self, message: str, __event_emitter__=None):
        if __event_emitter__:
            await __event_emitter__(
                {
                    "type": "notification",
                    "data": {"type": "error", "content": message},
                }
            )

    def _normalize_content(self, content: Any) -> str:
        if content is None:
            return ""
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            parts: list[str] = []
            for item in content:
                if not isinstance(item, dict):
                    continue
                if item.get("type") == "text":
                    parts.append(str(item.get("text", "")))
                elif isinstance(item.get("content"), str):
                    parts.append(item.get("content"))
            return "\n".join(part for part in parts if part)
        return str(content)

    def get_message_content(self, body: dict) -> str:
        target_id = body.get("id")
        messages = body.get("messages", []) or []

        if target_id:
            for message in messages:
                if not isinstance(message, dict):
                    continue

                candidate_ids = [
                    message.get("id"),
                    message.get("message_id"),
                    (
                        (message.get("info") or {}).get("id")
                        if isinstance(message.get("info"), dict)
                        else None
                    ),
                    (
                        (message.get("meta") or {}).get("id")
                        if isinstance(message.get("meta"), dict)
                        else None
                    ),
                ]

                if target_id in candidate_ids and message.get("role") == "assistant":
                    return self._normalize_content(message.get("content"))

        for message in reversed(messages):
            if isinstance(message, dict) and message.get("role") == "assistant":
                return self._normalize_content(message.get("content"))

        return ""

    def _extract_text(self, value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, str):
            return value
        if isinstance(value, list):
            parts: list[str] = []
            for item in value:
                if not isinstance(item, dict):
                    continue
                if item.get("type") == "text":
                    parts.append(str(item.get("text", "")))
                elif isinstance(item.get("content"), str):
                    parts.append(item.get("content"))
            return "\n".join(part for part in parts if part)
        return str(value)

    def _extract_chat_completion_text(self, response: Any) -> str:
        if response is None:
            return ""

        if isinstance(response, dict):
            choices = response.get("choices")
            if isinstance(choices, list) and choices:
                message = choices[0].get("message", {})
                if isinstance(message, dict):
                    return self._extract_text(message.get("content"))

            if "content" in response:
                return self._extract_text(response.get("content"))

        response_body = getattr(response, "body", None)
        if response_body:
            try:
                payload = json.loads(response_body.decode("utf-8"))
            except Exception:
                return ""
            return self._extract_chat_completion_text(payload)

        return ""

    def _get_user_model(self, user_data: Any) -> UserModel | None:
        if isinstance(user_data, UserModel):
            return user_data

        if not isinstance(user_data, dict) or not user_data:
            return None

        try:
            return UserModel(**user_data)
        except Exception:
            timestamp = int(time.time())
            fallback_user_data = {
                "id": user_data.get("id", ""),
                "email": user_data.get("email", ""),
                "username": user_data.get("username"),
                "role": user_data.get("role", "user"),
                "name": user_data.get("name") or user_data.get("email") or "User",
                "profile_image_url": user_data.get("profile_image_url"),
                "profile_banner_image_url": user_data.get(
                    "profile_banner_image_url"
                ),
                "bio": user_data.get("bio"),
                "gender": user_data.get("gender"),
                "date_of_birth": user_data.get("date_of_birth"),
                "timezone": user_data.get("timezone"),
                "presence_state": user_data.get("presence_state"),
                "status_emoji": user_data.get("status_emoji"),
                "status_message": user_data.get("status_message"),
                "status_expires_at": user_data.get("status_expires_at"),
                "info": user_data.get("info"),
                "settings": user_data.get("settings"),
                "oauth": user_data.get("oauth"),
                "scim": user_data.get("scim"),
                "last_active_at": user_data.get("last_active_at") or timestamp,
                "updated_at": user_data.get("updated_at") or timestamp,
                "created_at": user_data.get("created_at") or timestamp,
            }

            try:
                return UserModel(**fallback_user_data)
            except Exception:
                return None

    def _resolve_current_model_id(self, body: dict, __model__) -> str:
        return body.get("model") or (__model__ or {}).get("id") or ""

    def _strip_outer_code_fence(self, text: str) -> str:
        raw = (text or "").strip()
        if not raw:
            return ""

        match = WRAPPED_CODE_FENCE_RE.match(raw)
        if match:
            return match.group("body").strip()
        return raw

    def _derive_fallback_title(self, markdown_text: str, message_id: str) -> str:
        heading_match = HEADING_LINE_RE.search(markdown_text or "")
        if heading_match:
            return heading_match.group(1).strip()

        for line in (markdown_text or "").splitlines():
            cleaned = re.sub(r"^[#>*\-\d.\s]+", "", line).strip()
            if cleaned:
                return cleaned[:80]

        return f"Presentation {message_id}"

    def _ensure_document_title(self, markdown_text: str, fallback_title: str) -> str:
        text = (markdown_text or "").strip()
        if HEADING_LINE_RE.search(text):
            return text
        return f"# {fallback_title}\n\n{text}".strip()

    def build_extract_mermaid_png_js(self, message_id: str) -> str:
        scale = self.valves.mermaid_scale
        return f"""
const messageId = {json.dumps(message_id)};
const exportScale = {scale};

function findMessageElement(id) {{
  const selectors = [
    `#message-${{CSS.escape(id)}}`,
    `#${{CSS.escape(id)}}`,
    `[data-message-id="${{id}}"]`,
    `[data-id="${{id}}"]`,
    `.message-${{CSS.escape(id)}}`
  ];
  for (const selector of selectors) {{
    const el = document.querySelector(selector);
    if (el) return el;
  }}
  return null;
}}

function sleep(ms) {{
  return new Promise(resolve => setTimeout(resolve, ms));
}}

async function loadHtml2Canvas() {{
  if (typeof window.html2canvas === "function") return;

  await new Promise((resolve, reject) => {{
    const existing = document.querySelector('script[data-openwebui-mermaid-export="html2canvas-pro"]');
    if (existing) {{
      existing.addEventListener("load", () => resolve(), {{ once: true }});
      existing.addEventListener("error", () => reject(new Error("Failed to load html2canvas-pro")), {{ once: true }});
      return;
    }}

    const script = document.createElement("script");
    script.src = "https://cdn.jsdelivr.net/npm/html2canvas-pro@1.5.12/dist/html2canvas-pro.min.js";
    script.async = true;
    script.dataset.openwebuiMermaidExport = "html2canvas-pro";
    script.onload = () => resolve();
    script.onerror = () => reject(new Error("Failed to load html2canvas-pro"));
    document.head.appendChild(script);
  }});

  if (typeof window.html2canvas !== "function") {{
    throw new Error("html2canvas-pro loaded but window.html2canvas is unavailable.");
  }}
}}

async function waitForImages(root, timeoutMs = 10000) {{
  const images = Array.from(root.querySelectorAll("img"));
  if (!images.length) return;

  await Promise.race([
    Promise.all(images.map(img => {{
      if (img.complete) return Promise.resolve();
      return new Promise(resolve => {{
        const done = () => resolve();
        img.addEventListener("load", done, {{ once: true }});
        img.addEventListener("error", done, {{ once: true }});
      }});
    }})),
    sleep(timeoutMs)
  ]);
}}

const root = findMessageElement(messageId);
if (!root) {{
  throw new Error(`Could not find message element for message id: ${{messageId}}`);
}}

await loadHtml2Canvas();

const diagramElements = Array.from(
  root.querySelectorAll("svg[id^='mermaid-'], svg.flowchart")
);

const diagrams = [];

for (let index = 0; index < diagramElements.length; index++) {{
  const svg = diagramElements[index];

  const tempRoot = document.createElement("div");
  tempRoot.style.position = "fixed";
  tempRoot.style.left = "-100000px";
  tempRoot.style.top = "0";
  tempRoot.style.background = "#ffffff";
  tempRoot.style.padding = "0";
  tempRoot.style.margin = "0";
  tempRoot.style.zIndex = "-1";
  tempRoot.style.display = "inline-block";
  tempRoot.style.boxSizing = "border-box";

  const clone = svg.cloneNode(true);
  clone.removeAttribute("width");
  clone.removeAttribute("height");
  clone.style.display = "block";
  clone.style.margin = "0";
  clone.style.background = "#ffffff";

  const viewBox = svg.getAttribute("viewBox");
  if (viewBox) {{
    const parts = viewBox.trim().split(/\\s+/).map(Number);
    if (parts.length === 4 && Number.isFinite(parts[2]) && Number.isFinite(parts[3])) {{
      clone.setAttribute("width", String(parts[2]));
      clone.setAttribute("height", String(parts[3]));
      tempRoot.style.width = `${{parts[2]}}px`;
      tempRoot.style.height = `${{parts[3]}}px`;
    }}
  }}

  tempRoot.appendChild(clone);
  document.body.appendChild(tempRoot);

  try {{
    await waitForImages(tempRoot, 5000);
    await sleep(150);

    const canvas = await window.html2canvas(tempRoot, {{
      scale: exportScale,
      useCORS: true,
      backgroundColor: "#ffffff",
      logging: false,
      allowTaint: false,
      foreignObjectRendering: false
    }});

    const png = canvas.toDataURL("image/png");

    diagrams.push({{
      index,
      id: svg.id || null,
      png,
      width: canvas.width,
      height: canvas.height
    }});
  }} finally {{
    tempRoot.remove();
  }}
}}

return {{
  success: true,
  messageId,
  diagramCount: diagrams.length,
  diagrams
}};
"""

    def data_url_to_bytes(self, data_url: str) -> bytes:
        match = re.match(r"^data:[^;]+;base64,(.+)$", data_url, re.DOTALL)
        if not match:
            raise ValueError("Invalid data URL")
        return base64.b64decode(match.group(1))

    def fetch_image_bytes(self, src: str) -> bytes:
        src = (src or "").strip()
        if not src:
            raise ValueError("Empty image src")

        if src.startswith("data:"):
            self._debug_log("Loading image from data URL", src_preview=self._preview(src, 80))
            image_bytes = self.data_url_to_bytes(src)
            self._debug_log("Loaded image from data URL", size_bytes=len(image_bytes))
            return image_bytes

        parsed = urlparse(src)
        if parsed.scheme in {"http", "https"}:
            self._debug_log(
                "Downloading image",
                src=src,
                timeout_s=self.valves.request_timeout_s,
            )
            response = requests.get(src, timeout=self.valves.request_timeout_s)
            response.raise_for_status()
            self._debug_log(
                "Downloaded image",
                src=src,
                status_code=response.status_code,
                content_length=len(response.content),
            )
            return response.content

        raise ValueError(f"Unsupported image src: {src}")

    def get_image_info(self, src: str) -> tuple[bytes, int, int]:
        cached = self._image_info_cache.get(src)
        if cached is not None:
            self._debug_log(
                "Using cached image info",
                src=self._preview(src, 120),
                width_px=cached[1],
                height_px=cached[2],
                size_bytes=len(cached[0]),
            )
            return cached

        image_bytes = self.fetch_image_bytes(src)
        width = 0
        height = 0

        try:
            with Image.open(io.BytesIO(image_bytes)) as image:
                width, height = image.size
        except Exception:
            width = 0
            height = 0

        self._image_info_cache[src] = (image_bytes, width, height)
        self._debug_log(
            "Cached image info",
            src=self._preview(src, 120),
            width_px=width,
            height_px=height,
            size_bytes=len(image_bytes),
        )
        return image_bytes, width, height

    def _next_asset_token(self, assets: list[AssetToken]) -> str:
        return f"@@OPENWEBUI_ASSET_{len(assets) + 1}@@"

    def _replace_mermaid_with_assets(
        self,
        markdown_text: str,
        diagrams: list[dict],
    ) -> tuple[str, list[AssetToken]]:
        assets: list[AssetToken] = []
        diagram_iter = iter(diagrams or [])

        def repl(match: re.Match) -> str:
            code = match.group(1).strip()
            try:
                item = next(diagram_iter)
                png = str(item.get("png", ""))
                if not png.startswith("data:image/png;base64,"):
                    raise ValueError("Invalid Mermaid PNG payload")

                token = self._next_asset_token(assets)
                assets.append(
                    AssetToken(
                        token=token,
                        src=png,
                        alt="Mermaid diagram",
                        width_px=int(item.get("width", 0) or 0),
                        height_px=int(item.get("height", 0) or 0),
                    )
                )
                return f"\n\n{token}\n\n"
            except Exception:
                return f"\n```mermaid\n{code}\n```\n"

        return MERMAID_BLOCK_RE.sub(repl, markdown_text), assets

    def _replace_html_images_with_assets(
        self,
        markdown_text: str,
        assets: list[AssetToken],
    ) -> str:
        def repl(match: re.Match) -> str:
            tag_html = match.group(0)
            soup = BeautifulSoup(tag_html, "html.parser")
            img = soup.find("img")
            if not img:
                return tag_html

            src = (img.get("src") or "").strip()
            if not src:
                return tag_html

            token = self._next_asset_token(assets)
            assets.append(
                AssetToken(
                    token=token,
                    src=src,
                    alt=(img.get("alt") or "Image").strip() or "Image",
                )
            )
            return token

        return HTML_IMG_RE.sub(repl, markdown_text)

    def _replace_markdown_images_with_assets(
        self,
        markdown_text: str,
        assets: list[AssetToken],
    ) -> str:
        def repl(match: re.Match) -> str:
            raw_src = match.group("src").strip()
            src = raw_src.split()[0].strip() if raw_src else ""
            if src.startswith("<") and src.endswith(">"):
                src = src[1:-1]
            if not src:
                return match.group(0)

            token = self._next_asset_token(assets)
            assets.append(
                AssetToken(
                    token=token,
                    src=src,
                    alt=(match.group("alt") or "Image").strip() or "Image",
                )
            )
            return token

        return MARKDOWN_IMAGE_RE.sub(repl, markdown_text)

    def tokenize_assets(
        self,
        markdown_text: str,
        diagrams: list[dict],
    ) -> tuple[str, list[AssetToken]]:
        tokenized_text, assets = self._replace_mermaid_with_assets(
            markdown_text, diagrams
        )
        tokenized_text = self._replace_html_images_with_assets(tokenized_text, assets)
        tokenized_text = self._replace_markdown_images_with_assets(
            tokenized_text, assets
        )
        asset_type_counts = {
            "mermaid_or_image_assets": len(assets),
            "diagram_count": len(diagrams or []),
        }
        self._debug_log(
            "Tokenized visual assets from markdown",
            markdown_length=len(markdown_text),
            tokenized_length=len(tokenized_text),
            **asset_type_counts,
        )
        return tokenized_text, assets

    def _render_asset_html(self, asset: AssetToken) -> str:
        src = html.escape(asset.src, quote=True)
        alt = html.escape(asset.alt, quote=True)
        return f'<img src="{src}" alt="{alt}" />'

    def restore_assets(self, markdown_text: str, assets: list[AssetToken]) -> str:
        text = markdown_text
        missing_assets: list[AssetToken] = []

        for asset in assets:
            if asset.token in text:
                text = text.replace(asset.token, self._render_asset_html(asset))
            else:
                missing_assets.append(asset)

        if missing_assets:
            appendix = "\n\n".join(
                self._render_asset_html(asset) for asset in missing_assets
            )
            text = f"{text.rstrip()}\n\n{appendix}\n"

        self._debug_log(
            "Restored visual assets into markdown",
            markdown_length=len(markdown_text),
            restored_length=len(text),
            asset_count=len(assets),
            missing_asset_count=len(missing_assets),
        )
        return text

    async def structure_markdown_with_llm(
        self,
        markdown_text: str,
        body: dict,
        __request__,
        __user__,
        __model__,
        assets: list[AssetToken],
    ) -> str:
        if not self.valves.use_llm:
            self._debug_log("Skipping LLM restructuring because use_llm is disabled")
            return markdown_text

        if __request__ is None:
            self._debug_log(
                "Skipping LLM restructuring because __request__ is unavailable"
            )
            return markdown_text

        generation_model = self.valves.llm_model.strip() or self._resolve_current_model_id(
            body, __model__
        )
        if not generation_model:
            self._debug_log(
                "Skipping LLM restructuring because no model ID is available"
            )
            return markdown_text

        user_model = self._get_user_model(__user__)
        if user_model is None:
            self._debug_log(
                "Skipping LLM restructuring because UserModel could not be built"
            )
            return markdown_text

        asset_guidance = ""
        if assets:
            asset_lines = "\n".join(
                f"- {asset.token}: keep this token exactly unchanged"
                for asset in assets
            )
            asset_guidance = (
                "\n\nVisual tokens that must be preserved exactly:\n" f"{asset_lines}"
            )

        payload = {
            "model": generation_model,
            "stream": False,
            "temperature": self.valves.llm_temperature,
            "max_tokens": self.valves.llm_max_tokens,
            "messages": [
                {
                    "role": "system",
                    "content": self.valves.llm_prompt.strip() or DEFAULT_LLM_PROMPT,
                },
                {
                    "role": "user",
                    "content": (
                        "Rewrite the following assistant message into a well-structured "
                        "markdown document for presentation export.\n"
                        "Keep the content faithful. Use one H1 title.\n\n"
                        f"{markdown_text.strip()}{asset_guidance}"
                    ),
                },
            ],
        }

        self._debug_log(
            "Requesting structured markdown from LLM",
            generation_model=generation_model,
            temperature=self.valves.llm_temperature,
            max_tokens=self.valves.llm_max_tokens,
            markdown_length=len(markdown_text),
            asset_count=len(assets),
            prompt_preview=self._preview(self.valves.llm_prompt, 180),
        )

        try:
            response = await generate_chat_completion(
                __request__,
                payload,
                user_model,
                bypass_system_prompt=True,
            )
            structured = self._strip_outer_code_fence(
                self._extract_chat_completion_text(response)
            )
            self._debug_log(
                "Received structured markdown from LLM",
                structured_length=len(structured),
                used_fallback=not bool(structured),
                structured_preview=self._preview(structured, 220),
            )
            return structured or markdown_text
        except Exception as exc:
            self._debug_log(
                "LLM restructuring failed; falling back to original markdown",
                error=str(exc),
            )
            return markdown_text

    def _markdown_to_outline(
        self,
        markdown_text: str,
        fallback_title: str,
    ) -> OutlineNode:
        html_body = markdown.markdown(
            markdown_text,
            extensions=[
                "extra",
                "tables",
                "fenced_code",
                "sane_lists",
                "nl2br",
            ],
            output_format="html5",
        )

        soup = BeautifulSoup(html_body, "html.parser")
        container = soup.body if soup.body else soup

        root = OutlineNode(level=0, title="")
        stack: list[OutlineNode] = [root]

        for element in container.contents:
            if isinstance(element, NavigableString):
                text = str(element).strip()
                if text:
                    stack[-1].blocks.append(f"<p>{html.escape(text)}</p>")
                continue

            if not isinstance(element, Tag):
                continue

            tag = element.name.lower()
            if re.fullmatch(r"h[1-6]", tag):
                level = int(tag[1])
                title = element.get_text(" ", strip=True) or fallback_title

                if level == 1 and not root.title:
                    root.title = title
                    stack = [root]
                    continue

                node = OutlineNode(level=level, title=title)
                while stack and stack[-1].level >= level:
                    stack.pop()
                if not stack:
                    stack = [root]
                stack[-1].children.append(node)
                stack.append(node)
            else:
                stack[-1].blocks.append(str(element))

        root.title = root.title or fallback_title
        self._debug_log(
            "Built outline from markdown",
            outline_title=root.title,
            root_block_count=len(root.blocks),
            top_level_child_count=len(root.children),
            markdown_length=len(markdown_text),
        )
        return root

    def _html_block_to_content_blocks(self, block_html: str) -> list[ContentBlock]:
        soup = BeautifulSoup(block_html, "html.parser")
        container = soup.body if soup.body else soup
        blocks: list[ContentBlock] = []

        for element in container.contents:
            if isinstance(element, NavigableString):
                text = str(element).strip()
                if text:
                    blocks.append(ContentBlock(kind="paragraph", text=text))
                continue

            if not isinstance(element, Tag):
                continue

            tag = element.name.lower()

            if tag == "p":
                clone = BeautifulSoup(str(element), "html.parser")
                for image in clone.find_all("img"):
                    image.decompose()
                text = clone.get_text(" ", strip=True)
                if text:
                    blocks.append(ContentBlock(kind="paragraph", text=text))
                for image in element.find_all("img"):
                    src = (image.get("src") or "").strip()
                    if src:
                        blocks.append(
                            ContentBlock(
                                kind="image",
                                src=src,
                                alt=(image.get("alt") or "Image").strip() or "Image",
                            )
                        )

            elif tag in {"ul", "ol"}:
                items: list[str] = []
                for li in element.find_all("li", recursive=False):
                    li_clone = BeautifulSoup(str(li), "html.parser")
                    for nested in li_clone.find_all(["ul", "ol"]):
                        nested.decompose()
                    item_text = li_clone.get_text(" ", strip=True)
                    if item_text:
                        items.append(item_text)
                if items:
                    blocks.append(
                        ContentBlock(
                            kind="list",
                            items=items,
                            ordered=(tag == "ol"),
                        )
                    )

            elif tag == "blockquote":
                text = element.get_text(" ", strip=True)
                if text:
                    blocks.append(ContentBlock(kind="quote", text=text))

            elif tag == "pre":
                text = element.get_text("\n", strip=True)
                if text:
                    blocks.append(ContentBlock(kind="code", text=text))

            elif tag == "table":
                rows: list[list[str]] = []
                for tr in element.find_all("tr"):
                    row = [
                        cell.get_text(" ", strip=True)
                        for cell in tr.find_all(["th", "td"], recursive=False)
                    ]
                    if row:
                        rows.append(row)
                if rows:
                    blocks.append(ContentBlock(kind="table", rows=rows))

            elif tag == "img":
                src = (element.get("src") or "").strip()
                if src:
                    blocks.append(
                        ContentBlock(
                            kind="image",
                            src=src,
                            alt=(element.get("alt") or "Image").strip() or "Image",
                        )
                    )

            elif tag in {"div", "section", "article"}:
                nested_blocks: list[ContentBlock] = []
                for child in element.contents:
                    nested_blocks.extend(
                        self._html_block_to_content_blocks(str(child))
                    )
                if nested_blocks:
                    blocks.extend(nested_blocks)
                else:
                    text = element.get_text(" ", strip=True)
                    if text:
                        blocks.append(ContentBlock(kind="paragraph", text=text))

            elif tag == "hr":
                blocks.append(ContentBlock(kind="separator", text=""))

            else:
                text = element.get_text(" ", strip=True)
                if text:
                    blocks.append(ContentBlock(kind="paragraph", text=text))

        return blocks

    def _blocks_from_html_list(self, html_blocks: list[str]) -> list[ContentBlock]:
        blocks: list[ContentBlock] = []
        for block_html in html_blocks:
            blocks.extend(self._html_block_to_content_blocks(block_html))
        return blocks

    def _flatten_subtree_blocks(self, node: OutlineNode) -> list[ContentBlock]:
        blocks = self._blocks_from_html_list(node.blocks)
        for child in node.children:
            blocks.append(ContentBlock(kind="subheading", text=child.title))
            blocks.extend(self._flatten_subtree_blocks(child))
        return blocks

    def _estimate_block_weight(self, block: ContentBlock) -> float:
        if block.kind == "paragraph":
            return max(0.8, len(block.text) / 240.0)
        if block.kind == "quote":
            return max(0.9, len(block.text) / 220.0)
        if block.kind == "code":
            return max(1.1, len(block.text) / 300.0)
        if block.kind == "list":
            total_chars = sum(len(item) for item in block.items)
            return max(1.0, len(block.items) * 0.32 + total_chars / 260.0)
        if block.kind == "table":
            cols = max((len(row) for row in block.rows), default=1)
            return max(1.4, len(block.rows) * 0.35 + cols * 0.2)
        if block.kind == "image":
            return 1.8
        if block.kind == "subheading":
            return 0.45
        if block.kind == "separator":
            return 0.15
        return 1.0

    def _estimate_blocks_weight(self, blocks: list[ContentBlock]) -> float:
        return sum(self._estimate_block_weight(block) for block in blocks)

    def _split_text(self, text: str, target_chars: int) -> list[str]:
        stripped = (text or "").strip()
        if not stripped or len(stripped) <= target_chars:
            return [stripped] if stripped else []

        sentences = [
            part.strip() for part in SENTENCE_SPLIT_RE.split(stripped) if part.strip()
        ]
        if len(sentences) <= 1:
            words = stripped.split()
            chunks: list[str] = []
            current_words: list[str] = []
            current_len = 0
            for word in words:
                addition = len(word) + (1 if current_words else 0)
                if current_words and current_len + addition > target_chars:
                    chunks.append(" ".join(current_words))
                    current_words = [word]
                    current_len = len(word)
                else:
                    current_words.append(word)
                    current_len += addition
            if current_words:
                chunks.append(" ".join(current_words))
            return self._rebalance_chunks_to_avoid_trailing_colon(chunks)

        chunks: list[str] = []
        current_sentences: list[str] = []
        current_len = 0

        for sentence in sentences:
            addition = len(sentence) + (1 if current_sentences else 0)
            if current_sentences and current_len + addition > target_chars:
                chunks.append(" ".join(current_sentences))
                current_sentences = [sentence]
                current_len = len(sentence)
            else:
                current_sentences.append(sentence)
                current_len += addition

        if current_sentences:
            chunks.append(" ".join(current_sentences))

        return self._rebalance_chunks_to_avoid_trailing_colon(chunks)

    def _rebalance_chunks_to_avoid_trailing_colon(
        self,
        chunks: list[str],
    ) -> list[str]:
        rebalanced = [chunk.strip() for chunk in chunks if chunk and chunk.strip()]
        index = 0

        while index < len(rebalanced) - 1:
            current = rebalanced[index].rstrip()
            if not current.endswith(":"):
                index += 1
                continue

            split_index = max(
                current.rfind(". "),
                current.rfind("? "),
                current.rfind("! "),
                current.rfind("\n"),
            )

            if split_index >= 0:
                if current[split_index] == "\n":
                    head = current[:split_index].strip()
                    tail = current[split_index + 1 :].strip()
                else:
                    head = current[: split_index + 1].strip()
                    tail = current[split_index + 2 :].strip()

                if tail.endswith(":"):
                    if head:
                        rebalanced[index] = head
                    else:
                        del rebalanced[index]
                        index = max(0, index - 1)
                    rebalanced[index + 1 if head else index] = (
                        f"{tail} {rebalanced[index + 1 if head else index]}".strip()
                    )
                    index += 1
                    continue

            trailing_intro = current
            del rebalanced[index]
            rebalanced[index] = f"{trailing_intro} {rebalanced[index]}".strip()

        return [chunk for chunk in rebalanced if chunk.strip()]

    def _split_block_for_budget(
        self,
        block: ContentBlock,
        max_weight: float,
    ) -> list[ContentBlock]:
        if block.kind in {"paragraph", "quote", "code"}:
            if self._estimate_block_weight(block) <= max_weight:
                return [block]
            target_chars = max(140, int(max_weight * 220))
            chunks = self._split_text(block.text, target_chars)
            return [
                ContentBlock(kind=block.kind, text=chunk) for chunk in chunks if chunk
            ]

        if block.kind == "list":
            if self._estimate_block_weight(block) <= max_weight:
                return [block]

            grouped_items: list[list[str]] = []
            current_items: list[str] = []
            current_weight = 0.0
            for item in block.items:
                item_block = ContentBlock(
                    kind="list",
                    items=[item],
                    ordered=block.ordered,
                )
                item_weight = self._estimate_block_weight(item_block)
                if current_items and (
                    current_weight + item_weight > max_weight
                    or len(current_items) >= 5
                ):
                    grouped_items.append(current_items)
                    current_items = [item]
                    current_weight = item_weight
                else:
                    current_items.append(item)
                    current_weight += item_weight
            if current_items:
                grouped_items.append(current_items)

            return [
                ContentBlock(kind="list", items=items, ordered=block.ordered)
                for items in grouped_items
                if items
            ]

        if block.kind == "table" and len(block.rows) > 6:
            header = block.rows[:1]
            body_rows = block.rows[1:] if len(block.rows) > 1 else []
            chunks: list[ContentBlock] = []
            for index in range(0, len(body_rows), 5):
                rows = header + body_rows[index : index + 5]
                chunks.append(ContentBlock(kind="table", rows=rows))
            return chunks or [block]

        return [block]

    def _pack_text_blocks(
        self,
        blocks: list[ContentBlock],
        budget: float,
    ) -> list[list[ContentBlock]]:
        normalized: list[ContentBlock] = []
        split_budget = max(1.2, budget * 0.85)
        for block in blocks:
            normalized.extend(self._split_block_for_budget(block, split_budget))

        if not normalized:
            return [[]]

        chunks: list[list[ContentBlock]] = []
        current: list[ContentBlock] = []
        current_weight = 0.0

        for block in normalized:
            weight = self._estimate_block_weight(block)
            if current and current_weight + weight > budget:
                chunks.append(current)
                current = [block]
                current_weight = weight
            else:
                current.append(block)
                current_weight += weight

        if current:
            chunks.append(current)

        return chunks

    def _looks_like_bullet_grid(self, blocks: list[ContentBlock]) -> bool:
        if len(blocks) != 1:
            return False

        block = blocks[0]
        if block.kind != "list":
            return False

        if not (3 <= len(block.items) <= 5):
            return False

        return all(len(item) <= 90 for item in block.items)

    def _is_vertical_image(self, image_block: ContentBlock) -> bool:
        if not image_block.src:
            return False

        try:
            _, width, height = self.get_image_info(image_block.src)
        except Exception:
            return False

        return width > 0 and height > 0 and height > (width * 1.25)

    def _build_content_slide_specs(
        self,
        title: str,
        section_title: str,
        subsection_title: str,
        blocks: list[ContentBlock],
    ) -> list[SlideSpec]:
        blocks = self._normalize_render_blocks(blocks)
        if not blocks:
            return []

        images = [block for block in blocks if block.kind == "image"]
        text_blocks = [block for block in blocks if block.kind != "image"]

        def build_text_specs(
            chunk_budget: float,
            image: ContentBlock | None = None,
        ) -> list[SlideSpec]:
            chunked_blocks = self._pack_text_blocks(text_blocks, chunk_budget)
            return [
                SlideSpec(
                    kind="content",
                    title=title,
                    section_title=section_title,
                    subsection_title=subsection_title,
                    blocks=chunk,
                    image=image,
                )
                for chunk in chunked_blocks
            ]

        if len(images) == 1:
            image = images[0]
            if text_blocks:
                budget = 2.6 if self._is_vertical_image(image) else 3.0
                return build_text_specs(budget, image=image)

            return [
                SlideSpec(
                    kind="content",
                    title=title,
                    section_title=section_title,
                    subsection_title=subsection_title,
                    blocks=[],
                    image=image,
                )
            ]

        if not images:
            return build_text_specs(4.25, image=None)

        specs: list[SlideSpec] = []
        pending_text: list[ContentBlock] = []

        for block in blocks:
            if block.kind == "image":
                if pending_text:
                    specs.extend(
                        self._build_content_slide_specs(
                            title,
                            section_title,
                            subsection_title,
                            pending_text,
                        )
                    )
                    pending_text = []

                specs.append(
                    SlideSpec(
                        kind="content",
                        title=title,
                        section_title=section_title,
                        subsection_title=subsection_title,
                        blocks=[],
                        image=block,
                    )
                )
            else:
                pending_text.append(block)

        if pending_text:
            specs.extend(
                self._build_content_slide_specs(
                    title,
                    section_title,
                    subsection_title,
                    pending_text,
                )
            )

        return specs

    def _subtree_fits_single_slide_sequence(self, node: OutlineNode) -> bool:
        if not node.children:
            return False

        blocks = self._flatten_subtree_blocks(node)
        image_count = sum(1 for block in blocks if block.kind == "image")
        return image_count <= 1 and self._estimate_blocks_weight(blocks) <= 4.1

    def _plan_node_content(
        self,
        node: OutlineNode,
        section_title: str,
    ) -> list[SlideSpec]:
        specs: list[SlideSpec] = []

        if node.level >= 3 and self._subtree_fits_single_slide_sequence(node):
            blocks = self._flatten_subtree_blocks(node)
            specs.extend(
                self._build_content_slide_specs(
                    title=node.title,
                    section_title=section_title,
                    subsection_title=node.title,
                    blocks=blocks,
                )
            )
            return specs

        direct_blocks = self._blocks_from_html_list(node.blocks)
        if direct_blocks:
            specs.extend(
                self._build_content_slide_specs(
                    title=node.title,
                    section_title=section_title,
                    subsection_title=node.title if node.level >= 3 else "",
                    blocks=direct_blocks,
                )
            )

        for child in node.children:
            specs.extend(self._plan_node_content(child, section_title))

        return specs

    def _build_slide_plan(self, outline: OutlineNode) -> list[SlideSpec]:
        specs: list[SlideSpec] = [SlideSpec(kind="title", title=outline.title)]

        root_blocks = self._blocks_from_html_list(outline.blocks)
        if root_blocks:
            specs.extend(
                self._build_content_slide_specs(
                    title=outline.title,
                    section_title="",
                    subsection_title="",
                    blocks=root_blocks,
                )
            )

        level2_sections = [child for child in outline.children if child.level == 2]
        if level2_sections:
            for section in level2_sections:
                specs.append(
                    SlideSpec(
                        kind="section",
                        title=section.title,
                        section_title=section.title,
                    )
                )
                specs.extend(self._plan_node_content(section, section.title))

            for child in outline.children:
                if child.level != 2:
                    specs.extend(self._plan_node_content(child, ""))
        else:
            for child in outline.children:
                specs.extend(self._plan_node_content(child, ""))

        if self.valves.closing_slide_enabled:
            closing_blocks: list[ContentBlock] = []
            if self.valves.closing_slide_subtitle.strip():
                closing_blocks.append(
                    ContentBlock(
                        kind="paragraph",
                        text=self.valves.closing_slide_subtitle,
                    )
                )
            specs.append(
                SlideSpec(
                    kind="closing",
                    title=self.valves.closing_slide_title,
                    blocks=closing_blocks,
                )
            )

        return specs

    def _parse_color(self, value: str, default: str) -> RGBColor:
        raw = (value or "").strip().lstrip("#")
        if len(raw) == 3:
            raw = "".join(ch * 2 for ch in raw)
        if not re.fullmatch(r"[0-9A-Fa-f]{6}", raw):
            raw = default.strip().lstrip("#")
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))

    def _apply_placeholders(self, text: str, values: dict[str, str]) -> str:
        template = text or ""

        def repl(match: re.Match) -> str:
            key = match.group(1).strip().upper()
            return values.get(key, "")

        return PLACEHOLDER_RE.sub(repl, template)

    def _make_placeholder_values(
        self,
        slide_number: int,
        total_slides: int,
        file_name: str,
        presentation_title: str,
        section_title: str,
        subsection_title: str,
        user_name: str,
        export_now: datetime,
    ) -> dict[str, str]:
        export_date = export_now.strftime("%Y-%m-%d")
        export_time = export_now.strftime("%H:%M")
        export_timestamp = f"{export_date} {export_time}"

        values = {
            "CURRENT_PAGE": str(slide_number),
            "PAGE": str(slide_number),
            "TOTAL_PAGES": str(total_slides),
            "PAGES": str(total_slides),
            "FILE_NAME": file_name,
            "PRESENTATION_TITLE": presentation_title,
            "TITLE": presentation_title,
            "CURRENT_SECTION_TITLE": section_title,
            "SECTION_TITLE": section_title,
            "CURRENT_SUBSECTION_TITLE": subsection_title,
            "SUBSECTION_TITLE": subsection_title,
            "USER_NAME": user_name,
            "USERNAME": user_name,
            "EXPORT_DATE": export_date,
            "DATE": export_date,
            "EXPORT_TIME": export_time,
            "TIME": export_time,
            "EXPORT_TIMESTAMP": export_timestamp,
            "TIMESTAMP": export_timestamp,
            "HEADER_LOGO": "{{ HEADER_LOGO }}",
            "FOOTER_LOGO": "{{ FOOTER_LOGO }}",
        }

        values.update({key.lower(): value for key, value in values.items()})
        values.update({key.upper(): value for key, value in values.items()})
        return values

    def _add_textbox(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        text: str,
        font_name: str,
        font_size_pt: int,
        color_hex: str,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        bold: bool = False,
        italic: bool = False,
    ):
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_left = Pt(4)
        text_frame.margin_right = Pt(4)
        text_frame.margin_top = Pt(2)
        text_frame.margin_bottom = Pt(2)

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        paragraph.text = text or ""
        paragraph.space_after = Pt(0)

        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = self._parse_color(color_hex, "#000000")

        return shape

    def _apply_slide_background(self, slide, color_hex: str, image_url: str = ""):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._parse_color(color_hex, "#FFFFFF")

        if not image_url.strip():
            return

        try:
            image_bytes, _, _ = self.get_image_info(image_url.strip())
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                0,
                0,
                width=int(Inches(self.valves.slide_width_in)),
                height=int(Inches(self.valves.slide_height_in)),
            )
        except Exception:
            return

    def _split_bar_content_parts(
        self,
        template: str,
    ) -> list[tuple[str, str]]:
        value = template or ""
        if not value:
            return []

        parts: list[tuple[str, str]] = []
        cursor = 0
        for match in BAR_LOGO_PLACEHOLDER_RE.finditer(value):
            if match.start() > cursor:
                text_part = value[cursor : match.start()]
                if text_part:
                    parts.append(("text", text_part))
            parts.append(("logo", match.group(2).upper()))
            cursor = match.end()

        if cursor < len(value):
            text_part = value[cursor:]
            if text_part:
                parts.append(("text", text_part))

        return parts

    def _estimate_inline_text_width(self, text: str, font_size_pt: int) -> int:
        visible = text or ""
        if not visible:
            return 0
        estimated_points = max(8.0, (len(visible) * font_size_pt * 0.52) + 6.0)
        return int(Pt(estimated_points))

    def _measure_bar_logo(
        self,
        logo_key: str,
        logo_urls: dict[str, str],
        max_width: int,
        max_height: int,
    ) -> tuple[int, int]:
        logo_url = (logo_urls.get(logo_key) or "").strip()
        if not logo_url:
            return 0, 0

        try:
            _, image_width, image_height = self.get_image_info(logo_url)
        except Exception as exc:
            self._debug_log(
                "Failed to measure bar logo",
                logo_key=logo_key,
                logo_url=self._preview(logo_url, 120),
                error=str(exc),
            )
            return 0, 0

        fitted_width, fitted_height = self._fit_inside(
            image_width or max_width,
            image_height or max_height,
            max_width,
            max_height,
        )
        return fitted_width, fitted_height

    def _render_bar_logo(
        self,
        slide,
        logo_key: str,
        logo_urls: dict[str, str],
        left: int,
        top: int,
        width: int,
        height: int,
    ):
        logo_url = (logo_urls.get(logo_key) or "").strip()
        if not logo_url or width <= 0 or height <= 0:
            return

        try:
            image_bytes, image_width, image_height = self.get_image_info(logo_url)
            fitted_width, fitted_height = self._fit_inside(
                image_width or width,
                image_height or height,
                width,
                height,
            )
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                left + max(0, (width - fitted_width) // 2),
                top + max(0, (height - fitted_height) // 2),
                width=fitted_width,
                height=fitted_height,
            )
        except Exception as exc:
            self._debug_log(
                "Failed to render bar logo",
                logo_key=logo_key,
                logo_url=self._preview(logo_url, 120),
                error=str(exc),
            )

    def _render_bar_column_content(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        content: str,
        text_color: str,
        align: PP_ALIGN,
        logo_urls: dict[str, str],
    ):
        parts = self._split_bar_content_parts(content)
        if not parts:
            if content.strip():
                self._add_textbox(
                    slide,
                    left,
                    top,
                    width,
                    height,
                    content,
                    self.valves.body_font_name,
                    10,
                    text_color,
                    align=align,
                )
            return

        if all(part_kind == "text" for part_kind, _ in parts):
            self._add_textbox(
                slide,
                left,
                top,
                width,
                height,
                "".join(part_value for _, part_value in parts).strip(),
                self.valves.body_font_name,
                10,
                text_color,
                align=align,
            )
            return

        padding_x = int(Pt(4))
        padding_y = int(Pt(2))
        inner_left = left + padding_x
        inner_top = top + padding_y
        inner_width = max(1, width - (padding_x * 2))
        inner_height = max(1, height - (padding_y * 2))
        item_gap = int(Pt(4))

        measured_parts: list[tuple[str, str, int, int]] = []
        total_width = 0

        for part_kind, part_value in parts:
            if part_kind == "text":
                if not part_value.strip():
                    continue
                part_width = min(
                    inner_width,
                    self._estimate_inline_text_width(part_value, 10),
                )
                measured_parts.append((part_kind, part_value, part_width, inner_height))
                total_width += part_width
            else:
                logo_width, logo_height = self._measure_bar_logo(
                    part_value,
                    logo_urls,
                    inner_width,
                    inner_height,
                )
                if logo_width <= 0 or logo_height <= 0:
                    continue
                measured_parts.append((part_kind, part_value, logo_width, logo_height))
                total_width += logo_width

        if not measured_parts:
            return

        total_width += item_gap * max(0, len(measured_parts) - 1)
        if align == PP_ALIGN.CENTER:
            current_left = inner_left + max(0, (inner_width - total_width) // 2)
        elif align == PP_ALIGN.RIGHT:
            current_left = inner_left + max(0, inner_width - total_width)
        else:
            current_left = inner_left

        for index, (part_kind, part_value, part_width, part_height) in enumerate(
            measured_parts
        ):
            if part_kind == "text":
                self._add_textbox(
                    slide,
                    current_left,
                    inner_top,
                    part_width,
                    inner_height,
                    part_value.strip(),
                    self.valves.body_font_name,
                    10,
                    text_color,
                    align=PP_ALIGN.LEFT,
                )
            else:
                self._render_bar_logo(
                    slide,
                    part_value,
                    logo_urls,
                    current_left,
                    inner_top + max(0, (inner_height - part_height) // 2),
                    part_width,
                    part_height,
                )

            current_left += part_width
            if index < len(measured_parts) - 1:
                current_left += item_gap

    def _add_bar_with_three_columns(
        self,
        slide,
        top: int,
        height: int,
        bg_color: str,
        text_color: str,
        left_text: str,
        middle_text: str,
        right_text: str,
        logo_urls: dict[str, str],
    ):
        if height <= 0:
            return

        slide_width = int(Inches(self.valves.slide_width_in))
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, slide_width, height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._parse_color(bg_color, "#FFFFFF")
        bar.line.fill.background()

        left_width = int(slide_width * 0.44)
        middle_width = int(slide_width * 0.22)
        right_width = slide_width - left_width - middle_width
        self._render_bar_column_content(
            slide,
            0,
            top,
            left_width,
            height,
            left_text,
            text_color,
            align=PP_ALIGN.LEFT,
            logo_urls=logo_urls,
        )
        self._render_bar_column_content(
            slide,
            left_width,
            top,
            middle_width,
            height,
            middle_text,
            text_color,
            align=PP_ALIGN.CENTER,
            logo_urls=logo_urls,
        )
        self._render_bar_column_content(
            slide,
            left_width + middle_width,
            top,
            right_width,
            height,
            right_text,
            text_color,
            align=PP_ALIGN.RIGHT,
            logo_urls=logo_urls,
        )

    def _fit_inside(
        self,
        original_width: int,
        original_height: int,
        max_width: int,
        max_height: int,
    ) -> tuple[int, int]:
        if original_width <= 0 or original_height <= 0:
            return max_width, max_height

        ratio = min(max_width / original_width, max_height / original_height)
        ratio = max(ratio, 0.01)
        return (
            max(1, int(original_width * ratio)),
            max(1, int(original_height * ratio)),
        )

    def _get_fitted_image_dimensions(
        self,
        image_block: ContentBlock,
        width: int,
        height: int,
    ) -> tuple[int, int]:
        try:
            _, image_width, image_height = self.get_image_info(image_block.src)
        except Exception:
            image_width = width
            image_height = height

        return self._fit_inside(
            image_width or width,
            image_height or height,
            width,
            height,
        )

    def _normalize_render_blocks(
        self,
        blocks: list[ContentBlock],
    ) -> list[ContentBlock]:
        normalized: list[ContentBlock] = []

        for block in blocks:
            if block.kind in {"paragraph", "quote", "code", "subheading"}:
                if not block.text.strip():
                    continue
                normalized.append(block)
                continue

            if block.kind == "list":
                items = [item.strip() for item in block.items if item and item.strip()]
                if not items:
                    continue
                normalized.append(
                    ContentBlock(
                        kind=block.kind,
                        items=items,
                        ordered=block.ordered,
                    )
                )
                continue

            if block.kind == "table":
                rows = [
                    [cell.strip() for cell in row]
                    for row in block.rows
                    if any((cell or "").strip() for cell in row)
                ]
                if not rows:
                    continue
                normalized.append(ContentBlock(kind="table", rows=rows))
                continue

            if block.kind == "separator":
                if normalized and normalized[-1].kind != "separator":
                    normalized.append(block)
                continue

            normalized.append(block)

        while normalized and normalized[-1].kind == "separator":
            normalized.pop()

        return normalized

    def _estimate_chars_per_line(self, width: int, font_size_pt: int) -> int:
        width_in = max(1.0, width / 914400.0)
        scale = 16.0 / max(font_size_pt, 10)
        return max(12, int(width_in * 12.0 * scale))

    def _estimate_line_count(
        self,
        text: str,
        width: int,
        font_size_pt: int,
        prefix_chars: int = 0,
    ) -> int:
        chars_per_line = max(
            8,
            self._estimate_chars_per_line(width, font_size_pt) - prefix_chars,
        )
        return max(1, math.ceil(len((text or "").strip()) / chars_per_line))

    def _estimate_block_height(self, block: ContentBlock, width: int) -> int:
        if block.kind == "paragraph":
            lines = self._estimate_line_count(
                block.text,
                width,
                self.valves.body_font_size_pt,
            )
            return int(Pt(lines * self.valves.body_font_size_pt * 1.4 + 6))

        if block.kind == "quote":
            lines = self._estimate_line_count(
                block.text,
                width,
                self.valves.body_font_size_pt,
            )
            return int(Pt(lines * self.valves.body_font_size_pt * 1.45 + 8))

        if block.kind == "code":
            lines = max(1, len((block.text or "").splitlines()))
            return int(Pt(lines * max(11, self.valves.body_font_size_pt - 2) * 1.3 + 8))

        if block.kind == "subheading":
            font_size = max(14, self.valves.body_font_size_pt + 1)
            lines = self._estimate_line_count(block.text, width, font_size)
            return int(Pt(lines * font_size * 1.25 + 4))

        if block.kind == "list":
            total_lines = 0
            for index, item in enumerate(block.items, start=1):
                prefix_chars = len(f"{index}. ") if block.ordered else 2
                total_lines += self._estimate_line_count(
                    item,
                    width,
                    self.valves.body_font_size_pt,
                    prefix_chars=prefix_chars,
                )
            return int(Pt(total_lines * self.valves.body_font_size_pt * 1.35 + 6))

        if block.kind == "table":
            row_count = max(1, len(block.rows))
            return int(Pt(row_count * (self.valves.body_font_size_pt * 1.55) + 10))

        if block.kind == "separator":
            return int(Pt(6))

        return int(Pt(self.valves.body_font_size_pt * 1.4 + 6))

    def _estimate_blocks_stack_height(self, blocks: list[ContentBlock], width: int) -> int:
        render_blocks = self._normalize_render_blocks(blocks)
        if not render_blocks:
            return 0

        gap = int(Pt(6))
        total = sum(self._estimate_block_height(block, width) for block in render_blocks)
        total += gap * max(0, len(render_blocks) - 1)
        return total

    def _set_table_cell_text(
        self,
        cell,
        text: str,
        *,
        bold: bool = False,
        fill_color: str | None = None,
        font_color: str | None = None,
    ):
        if fill_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._parse_color(fill_color, "#FFFFFF")

        text_frame = cell.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_left = Pt(4)
        text_frame.margin_right = Pt(4)
        text_frame.margin_top = Pt(2)
        text_frame.margin_bottom = Pt(2)

        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = PP_ALIGN.LEFT

        for run in paragraph.runs:
            run.font.name = self.valves.body_font_name
            run.font.size = Pt(max(12, self.valves.body_font_size_pt - 1))
            run.font.bold = bold
            run.font.color.rgb = self._parse_color(
                font_color or self.valves.body_text_color,
                "#1E293B",
            )

    def _render_table_block(
        self,
        slide,
        rows: list[list[str]],
        left: int,
        top: int,
        width: int,
        height: int,
    ):
        row_count = len(rows)
        col_count = max((len(row) for row in rows), default=1)
        table = slide.shapes.add_table(
            row_count,
            col_count,
            left,
            top,
            width,
            height,
        ).table

        base_col_width = width // col_count
        for index in range(col_count):
            table.columns[index].width = (
                width - (base_col_width * index)
                if index == col_count - 1
                else base_col_width
            )

        base_row_height = max(1, height // row_count)
        for index in range(row_count):
            table.rows[index].height = (
                height - (base_row_height * index)
                if index == row_count - 1
                else base_row_height
            )

        for row_index, row in enumerate(rows):
            is_header = row_index == 0
            for col_index in range(col_count):
                cell = table.cell(row_index, col_index)
                text = row[col_index] if col_index < len(row) else ""
                self._set_table_cell_text(
                    cell,
                    text,
                    bold=is_header,
                    fill_color="#E2E8F0" if is_header else "#FFFFFF",
                    font_color="#0F172A" if is_header else self.valves.body_text_color,
                )

    def _render_block_stack(
        self,
        slide,
        blocks: list[ContentBlock],
        left: int,
        top: int,
        width: int,
        max_height: int,
    ) -> int:
        render_blocks = self._normalize_render_blocks(blocks)
        if not render_blocks:
            return 0

        gap = int(Pt(6))
        heights = [self._estimate_block_height(block, width) for block in render_blocks]
        total_height = sum(heights) + gap * max(0, len(render_blocks) - 1)
        current_top = top + max(0, (max_height - total_height) // 2)

        for index, (block, block_height) in enumerate(zip(render_blocks, heights)):
            remaining_height = top + max_height - current_top
            draw_height = min(block_height, max(1, remaining_height))

            if block.kind == "table":
                self._render_table_block(
                    slide,
                    block.rows,
                    left,
                    current_top,
                    width,
                    draw_height,
                )
            else:
                self._render_blocks_in_text_box(
                    slide,
                    [block],
                    left,
                    current_top,
                    width,
                    draw_height,
                )

            current_top += draw_height
            if index < len(render_blocks) - 1:
                current_top += gap

        return min(total_height, max_height)

    def _add_image_in_box(
        self,
        slide,
        image_block: ContentBlock,
        left: int,
        top: int,
        width: int,
        height: int,
    ):
        try:
            image_bytes, image_width, image_height = self.get_image_info(image_block.src)
            fitted_width, fitted_height = self._fit_inside(
                image_width or width,
                image_height or height,
                width,
                height,
            )
            fitted_left = left + max(0, (width - fitted_width) // 2)
            fitted_top = top + max(0, (height - fitted_height) // 2)
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                fitted_left,
                fitted_top,
                width=fitted_width,
                height=fitted_height,
            )
        except Exception:
            self._add_textbox(
                slide,
                left,
                top,
                width,
                height,
                f"[Image could not be loaded: {image_block.alt or image_block.src}]",
                self.valves.body_font_name,
                self.valves.body_font_size_pt,
                self.valves.body_text_color,
            )

    def _append_paragraph(
        self,
        text_frame,
        text: str,
        font_name: str,
        font_size_pt: int,
        color_hex: str,
        bold: bool = False,
        italic: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ):
        if not text.strip():
            return

        if text_frame.text:
            paragraph = text_frame.add_paragraph()
        else:
            paragraph = text_frame.paragraphs[0]

        paragraph.text = text.strip()
        paragraph.alignment = align
        paragraph.space_after = Pt(6)

        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = self._parse_color(color_hex, "#000000")

    def _render_table_text(self, rows: list[list[str]]) -> list[str]:
        return [" | ".join(cell for cell in row if cell) for row in rows if row]

    def _render_blocks_in_text_box(
        self,
        slide,
        blocks: list[ContentBlock],
        left: int,
        top: int,
        width: int,
        height: int,
    ):
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.margin_left = Pt(4)
        text_frame.margin_right = Pt(4)
        text_frame.margin_top = Pt(2)
        text_frame.margin_bottom = Pt(2)

        for block in blocks:
            if block.kind == "paragraph":
                self._append_paragraph(
                    text_frame,
                    block.text,
                    self.valves.body_font_name,
                    self.valves.body_font_size_pt,
                    self.valves.body_text_color,
                )
            elif block.kind == "quote":
                self._append_paragraph(
                    text_frame,
                    f'"{block.text}"',
                    self.valves.body_font_name,
                    self.valves.body_font_size_pt,
                    self.valves.body_text_color,
                    italic=True,
                )
            elif block.kind == "code":
                self._append_paragraph(
                    text_frame,
                    block.text,
                    "Consolas",
                    max(11, self.valves.body_font_size_pt - 2),
                    self.valves.body_text_color,
                )
            elif block.kind == "list":
                for index, item in enumerate(block.items, start=1):
                    prefix = f"{index}. " if block.ordered else "• "
                    self._append_paragraph(
                        text_frame,
                        f"{prefix}{item}",
                        self.valves.body_font_name,
                        self.valves.body_font_size_pt,
                        self.valves.body_text_color,
                    )
            elif block.kind == "table":
                for line in self._render_table_text(block.rows):
                    self._append_paragraph(
                        text_frame,
                        line,
                        self.valves.body_font_name,
                        max(12, self.valves.body_font_size_pt - 2),
                        self.valves.body_text_color,
                    )
            elif block.kind == "subheading":
                self._append_paragraph(
                    text_frame,
                    block.text,
                    self.valves.body_title_font_name,
                    max(14, self.valves.body_font_size_pt + 1),
                    self.valves.accent_color,
                    bold=True,
                )
            elif block.kind == "separator":
                self._append_paragraph(
                    text_frame,
                    " ",
                    self.valves.body_font_name,
                    4,
                    self.valves.body_text_color,
                )

    def _render_bullet_grid(
        self,
        slide,
        items: list[str],
        left: int,
        top: int,
        width: int,
        height: int,
    ):
        count = len(items)
        columns = 3 if count >= 5 else 2
        rows = (count + columns - 1) // columns
        gap = int(Inches(0.16))
        card_width = max(1, (width - (gap * (columns - 1))) // columns)
        max_card_height = max(1, (height - (gap * (rows - 1))) // rows)
        card_height = min(max_card_height, int(Inches(1.45)))
        total_grid_height = (card_height * rows) + (gap * max(0, rows - 1))
        start_top = top + max(0, (height - total_grid_height) // 2)

        for index, item in enumerate(items):
            row = index // columns
            col = index % columns
            card_left = left + col * (card_width + gap)
            card_top = start_top + row * (card_height + gap)

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                card_left,
                card_top,
                card_width,
                card_height,
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self._parse_color(
                self.valves.bullet_grid_bg_color,
                "#DBEAFE",
            )
            card.line.color.rgb = self._parse_color(self.valves.accent_color, "#2563EB")

            self._add_textbox(
                slide,
                card_left + int(Inches(0.06)),
                card_top + int(Inches(0.04)),
                card_width - int(Inches(0.12)),
                card_height - int(Inches(0.08)),
                item,
                self.valves.body_font_name,
                self.valves.body_font_size_pt,
                self.valves.bullet_grid_text_color,
                align=PP_ALIGN.CENTER,
                bold=True,
            )

    def _render_content_slide(
        self,
        presentation: Presentation,
        spec: SlideSpec,
        placeholder_values: dict[str, str],
    ):
        render_blocks = self._normalize_render_blocks(spec.blocks)
        has_renderable_text = bool(render_blocks)
        self._debug_log(
            "Rendering content slide",
            title=spec.title,
            section_title=spec.section_title,
            subsection_title=spec.subsection_title,
            block_count=len(render_blocks),
            has_image=bool(spec.image),
            bullet_grid=self._looks_like_bullet_grid(render_blocks),
        )
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._apply_slide_background(slide, self.valves.content_slide_bg_color)

        slide_width = int(Inches(self.valves.slide_width_in))
        slide_height = int(Inches(self.valves.slide_height_in))
        header_height = int(Inches(max(0, self.valves.header_height_in)))
        footer_height = int(Inches(max(0, self.valves.footer_height_in)))
        margin = int(Inches(self.valves.content_margin_in))

        content_left = margin
        content_top = header_height + margin
        content_width = slide_width - (margin * 2)
        content_height = slide_height - header_height - footer_height - (margin * 2)

        title_height = 0
        if spec.title.strip():
            title_height = int(Inches(0.7))
            self._add_textbox(
                slide,
                content_left,
                content_top,
                content_width,
                title_height,
                self._apply_placeholders(spec.title, placeholder_values),
                self.valves.body_title_font_name,
                self.valves.body_title_font_size_pt,
                self.valves.body_title_color,
                bold=True,
            )
            content_top += title_height + int(Inches(0.08))
            content_height -= title_height + int(Inches(0.08))

        if spec.image and has_renderable_text:
            if self._is_vertical_image(spec.image):
                gap = int(Inches(0.24))
                text_width = int(content_width * 0.58)
                image_width = content_width - text_width - gap
                text_height = min(
                    content_height,
                    self._estimate_blocks_stack_height(render_blocks, text_width),
                )
                _, fitted_image_height = self._get_fitted_image_dimensions(
                    spec.image,
                    image_width,
                    content_height,
                )
                combined_height = max(text_height, fitted_image_height)
                block_top = content_top + max(0, (content_height - combined_height) // 2)

                self._render_block_stack(
                    slide,
                    render_blocks,
                    content_left,
                    block_top + max(0, (combined_height - text_height) // 2),
                    text_width,
                    text_height,
                )
                self._add_image_in_box(
                    slide,
                    spec.image,
                    content_left + text_width + gap,
                    block_top + max(0, (combined_height - fitted_image_height) // 2),
                    image_width,
                    fitted_image_height,
                )
            else:
                gap = int(Inches(0.16))
                estimated_text_height = self._estimate_blocks_stack_height(
                    render_blocks,
                    content_width,
                )
                max_text_height = max(1, content_height - int(Inches(1.2)))
                text_height = min(estimated_text_height, max_text_height)
                image_available_height = max(1, content_height - text_height - gap)
                _, fitted_image_height = self._get_fitted_image_dimensions(
                    spec.image,
                    content_width,
                    image_available_height,
                )
                combined_height = text_height + gap + fitted_image_height
                block_top = content_top + max(0, (content_height - combined_height) // 2)

                self._render_block_stack(
                    slide,
                    render_blocks,
                    content_left,
                    block_top,
                    content_width,
                    text_height,
                )
                self._add_image_in_box(
                    slide,
                    spec.image,
                    content_left,
                    block_top + text_height + gap,
                    content_width,
                    fitted_image_height,
                )
        elif spec.image:
            _, fitted_image_height = self._get_fitted_image_dimensions(
                spec.image,
                content_width,
                content_height,
            )
            self._add_image_in_box(
                slide,
                spec.image,
                content_left,
                content_top + max(0, (content_height - fitted_image_height) // 2),
                content_width,
                fitted_image_height,
            )
        elif self._looks_like_bullet_grid(render_blocks):
            self._render_bullet_grid(
                slide,
                render_blocks[0].items,
                content_left,
                content_top,
                content_width,
                content_height,
            )
        else:
            self._render_block_stack(
                slide,
                render_blocks,
                content_left,
                content_top,
                content_width,
                content_height,
            )

        self._add_slide_header_footer(slide, placeholder_values)

    def _render_closing_slide(
        self,
        presentation: Presentation,
        spec: SlideSpec,
        placeholder_values: dict[str, str],
    ):
        self._debug_log("Rendering closing slide", title=spec.title)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._apply_slide_background(
            slide,
            self.valves.closing_slide_bg_color,
            self.valves.closing_slide_bg_image_url,
        )

        slide_width = int(Inches(self.valves.slide_width_in))
        slide_height = int(Inches(self.valves.slide_height_in))
        title_left = int(Inches(1.0))
        title_top = int(slide_height * 0.28)
        title_width = slide_width - int(Inches(2.0))
        title_height = int(Inches(1.4))
        subtitle_text = (
            self._apply_placeholders(spec.blocks[0].text, placeholder_values)
            if spec.blocks
            else ""
        )

        self._add_textbox(
            slide,
            title_left,
            title_top,
            title_width,
            title_height,
            self._apply_placeholders(spec.title, placeholder_values),
            self.valves.title_font_name,
            self.valves.title_font_size_pt,
            self.valves.closing_slide_title_color,
            align=PP_ALIGN.CENTER,
            bold=True,
        )
        if subtitle_text:
            self._add_textbox(
                slide,
                title_left,
                title_top + int(Inches(1.35)),
                title_width,
                int(Inches(0.9)),
                subtitle_text,
                self.valves.subtitle_font_name,
                self.valves.subtitle_font_size_pt,
                self.valves.closing_slide_subtitle_color,
                align=PP_ALIGN.CENTER,
            )

        self._add_slide_header_footer(slide, placeholder_values)

    def _render_title_slide(
        self,
        presentation: Presentation,
        spec: SlideSpec,
        placeholder_values: dict[str, str],
    ):
        self._debug_log("Rendering title slide", title=spec.title)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._apply_slide_background(
            slide,
            self.valves.title_slide_bg_color,
            self.valves.title_slide_bg_image_url,
        )

        slide_width = int(Inches(self.valves.slide_width_in))
        slide_height = int(Inches(self.valves.slide_height_in))
        title_left = int(Inches(1.0))
        title_top = int(slide_height * 0.24)
        title_width = slide_width - int(Inches(2.0))
        title_height = int(Inches(1.6))

        subtitle_text = self._apply_placeholders(
            self.valves.title_slide_subtitle,
            placeholder_values,
        )

        self._add_textbox(
            slide,
            title_left,
            title_top,
            title_width,
            title_height,
            self._apply_placeholders(spec.title, placeholder_values),
            self.valves.title_font_name,
            self.valves.title_font_size_pt,
            self.valves.title_color,
            align=PP_ALIGN.CENTER,
            bold=True,
        )
        self._add_textbox(
            slide,
            title_left,
            title_top + int(Inches(1.55)),
            title_width,
            int(Inches(0.8)),
            subtitle_text,
            self.valves.subtitle_font_name,
            self.valves.subtitle_font_size_pt,
            self.valves.subtitle_color,
            align=PP_ALIGN.CENTER,
        )

        self._add_slide_header_footer(slide, placeholder_values)

    def _render_section_slide(
        self,
        presentation: Presentation,
        spec: SlideSpec,
        placeholder_values: dict[str, str],
    ):
        self._debug_log(
            "Rendering section slide",
            title=spec.title,
            section_title=spec.section_title,
        )
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._apply_slide_background(
            slide,
            self.valves.section_title_slide_bg_color,
            self.valves.section_title_slide_bg_image_url,
        )

        slide_width = int(Inches(self.valves.slide_width_in))
        slide_height = int(Inches(self.valves.slide_height_in))
        title_left = int(Inches(0.9))
        title_top = int(slide_height * 0.33)
        title_width = slide_width - int(Inches(1.8))
        title_height = int(Inches(1.2))

        self._add_textbox(
            slide,
            title_left,
            title_top,
            title_width,
            title_height,
            self._apply_placeholders(spec.title, placeholder_values),
            self.valves.section_title_font_name,
            self.valves.section_title_font_size_pt,
            self.valves.section_title_color,
            align=PP_ALIGN.CENTER,
            bold=True,
        )

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(slide_width * 0.3),
            title_top + int(Inches(1.0)),
            int(slide_width * 0.4),
            int(Inches(0.06)),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self._parse_color(
            self.valves.accent_color,
            "#2563EB",
        )
        accent.line.fill.background()

        self._add_slide_header_footer(slide, placeholder_values)

    def _add_slide_header_footer(self, slide, placeholder_values: dict[str, str]):
        header_height = int(Inches(max(0, self.valves.header_height_in)))
        footer_height = int(Inches(max(0, self.valves.footer_height_in)))
        slide_height = int(Inches(self.valves.slide_height_in))
        logo_urls = {
            "HEADER_LOGO": self.valves.header_logo_url.strip(),
            "FOOTER_LOGO": self.valves.footer_logo_url.strip(),
        }

        if header_height > 0:
            self._add_bar_with_three_columns(
                slide,
                0,
                header_height,
                self.valves.header_bg_color,
                self.valves.header_text_color,
                self._apply_placeholders(
                    self.valves.header_text_left,
                    placeholder_values,
                ),
                self._apply_placeholders(
                    self.valves.header_text_middle,
                    placeholder_values,
                ),
                self._apply_placeholders(
                    self.valves.header_text_right,
                    placeholder_values,
                ),
                logo_urls=logo_urls,
            )

        if footer_height > 0:
            self._add_bar_with_three_columns(
                slide,
                slide_height - footer_height,
                footer_height,
                self.valves.footer_bg_color,
                self.valves.footer_text_color,
                self._apply_placeholders(
                    self.valves.footer_text_left,
                    placeholder_values,
                ),
                self._apply_placeholders(
                    self.valves.footer_text_middle,
                    placeholder_values,
                ),
                self._apply_placeholders(
                    self.valves.footer_text_right,
                    placeholder_values,
                ),
                logo_urls=logo_urls,
            )

    def build_pptx(
        self,
        markdown_text: str,
        file_name: str,
        user_name: str,
    ) -> tuple[bytes, int, str]:
        fallback_title = self._derive_fallback_title(markdown_text, file_name)
        outlined_markdown = self._ensure_document_title(markdown_text, fallback_title)
        outline = self._markdown_to_outline(outlined_markdown, fallback_title)
        slide_specs = self._build_slide_plan(outline)
        self._debug_log(
            "Planned slide deck",
            file_name=file_name,
            presentation_title=outline.title,
            markdown_length=len(markdown_text),
            slide_count=len(slide_specs),
            slide_summary=self._summarize_slide_specs(slide_specs),
        )

        presentation = Presentation()
        presentation.slide_width = Inches(self.valves.slide_width_in)
        presentation.slide_height = Inches(self.valves.slide_height_in)

        total_slides = len(slide_specs)
        export_now = datetime.now()

        for slide_number, spec in enumerate(slide_specs, start=1):
            placeholder_values = self._make_placeholder_values(
                slide_number=slide_number,
                total_slides=total_slides,
                file_name=file_name,
                presentation_title=outline.title,
                section_title=spec.section_title,
                subsection_title=spec.subsection_title,
                user_name=user_name,
                export_now=export_now,
            )

            if spec.kind == "title":
                self._render_title_slide(presentation, spec, placeholder_values)
            elif spec.kind == "section":
                self._render_section_slide(presentation, spec, placeholder_values)
            elif spec.kind == "closing":
                self._render_closing_slide(presentation, spec, placeholder_values)
            else:
                self._render_content_slide(presentation, spec, placeholder_values)

        output = io.BytesIO()
        presentation.save(output)
        self._debug_log(
            "Built PPTX bytes",
            file_name=file_name,
            slide_count=total_slides,
            size_bytes=output.tell(),
        )
        return output.getvalue(), total_slides, outline.title

    async def download_file(
        self,
        pptx_bytes: bytes,
        filename: str,
        __event_emitter__=None,
        __event_call__=None,
    ):
        encoded = base64.b64encode(pptx_bytes).decode("ascii")
        self._debug_log(
            "Preparing browser download",
            filename=filename,
            size_bytes=len(pptx_bytes),
            encoded_length=len(encoded),
            use_event_call=__event_call__ is not None,
            use_event_emitter=__event_emitter__ is not None,
        )

        js_code = f"""
const base64 = {json.dumps(encoded)};
const filename = {json.dumps(filename)};
const binary = atob(base64);
const bytes = new Uint8Array(binary.length);

for (let i = 0; i < binary.length; i++) {{
  bytes[i] = binary.charCodeAt(i);
}}

const blob = new Blob(
  [bytes],
  {{ type: "application/vnd.openxmlformats-officedocument.presentationml.presentation" }}
);
const url = URL.createObjectURL(blob);

try {{
  const a = document.createElement("a");
  a.href = url;
  a.download = filename;
  a.style.display = "none";
  document.body.appendChild(a);
  a.click();
  a.remove();
}} finally {{
  setTimeout(() => URL.revokeObjectURL(url), 4000);
}}

return {{ success: true, filename, size: bytes.length }};
"""

        payload = {"type": "execute", "data": {"code": js_code}}

        if __event_call__ is not None:
            result = await __event_call__(payload)
            self._debug_log(
                "Browser download execute completed via __event_call__",
                filename=filename,
                result=result,
            )
            return result

        if __event_emitter__ is not None:
            await __event_emitter__(payload)
            self._debug_log(
                "Browser download execute emitted via __event_emitter__",
                filename=filename,
                size_bytes=len(pptx_bytes),
            )
            return {"success": True, "filename": filename, "size": len(pptx_bytes)}

        self._debug_log(
            "Browser download skipped because no event channel is available",
            filename=filename,
        )
        return None

    async def action(
        self,
        body: dict,
        __user__=None,
        __event_emitter__=None,
        __event_call__=None,
        __request__=None,
        __model__=None,
        **kwargs,
    ):
        self._debug_log(
            "Action invoked",
            body_keys=sorted(body.keys()),
            chat_id=body.get("chat_id"),
            message_id=body.get("id"),
            model=body.get("model"),
            user_id=__user__.get("id") if isinstance(__user__, dict) else None,
            use_llm=self.valves.use_llm,
            llm_model_override=self.valves.llm_model,
        )
        message_id = body.get("id")
        if not message_id:
            await self.emit_error(
                "PPTX export failed: could not determine the current message id.",
                __event_emitter__,
            )
            return {
                "content": "Could not determine the current message id from body['id']."
            }

        filename = self.build_filename(message_id)

        await self.emit_status("Reading message content...", False, __event_emitter__)
        markdown_text = self.get_message_content(body)
        self._debug_log(
            "Extracted assistant message content",
            message_id=message_id,
            markdown_length=len(markdown_text),
            markdown_preview=self._preview(markdown_text, 220),
        )
        if not markdown_text.strip():
            await self.emit_error(
                "PPTX export failed: no assistant message content found.",
                __event_emitter__,
            )
            await self.emit_status(
                "No assistant message content found.",
                True,
                __event_emitter__,
            )
            return {"content": "No assistant message content found."}

        await self.emit_status(
            "Collecting rendered Mermaid diagrams...",
            False,
            __event_emitter__,
        )

        diagrams = []
        extract_result = None
        extract_error = None

        try:
            if __event_call__ is None:
                raise RuntimeError(
                    "This action needs __event_call__ so browser JS can return Mermaid image data."
                )

            extract_result = await __event_call__(
                {
                    "type": "execute",
                    "data": {"code": self.build_extract_mermaid_png_js(message_id)},
                }
            )

            if not isinstance(extract_result, dict):
                raise RuntimeError(
                    f"Unexpected execute result type: {type(extract_result).__name__}"
                )

            diagrams = extract_result.get("diagrams", []) or []
            self._debug_log(
                "Collected Mermaid diagrams",
                diagram_count=len(diagrams),
                extract_result_keys=sorted(extract_result.keys()),
            )
        except Exception as exc:
            extract_error = str(exc)
            diagrams = []
            self._debug_log(
                "Mermaid extraction failed; continuing without embedded diagrams",
                error=extract_error,
            )

        tokenized_markdown, assets = self.tokenize_assets(markdown_text, diagrams)
        self._debug_log(
            "Prepared tokenized markdown",
            asset_count=len(assets),
            tokenized_length=len(tokenized_markdown),
        )

        await self.emit_status(
            "Structuring the document for presentation export...",
            False,
            __event_emitter__,
        )

        fallback_title = self._derive_fallback_title(markdown_text, message_id)
        structured_markdown = await self.structure_markdown_with_llm(
            tokenized_markdown,
            body,
            __request__,
            __user__,
            __model__,
            assets,
        )
        structured_markdown = self._ensure_document_title(
            structured_markdown,
            fallback_title,
        )
        rendered_markdown = self.restore_assets(structured_markdown, assets)
        self._debug_log(
            "Prepared final markdown for PPTX generation",
            fallback_title=fallback_title,
            structured_length=len(structured_markdown),
            rendered_length=len(rendered_markdown),
            rendered_preview=self._preview(rendered_markdown, 220),
        )

        await self.emit_status("Generating PPTX...", False, __event_emitter__)

        user_name = ""
        if isinstance(__user__, dict):
            user_name = (__user__.get("name") or __user__.get("email") or "").strip()

        try:
            pptx_bytes, slide_count, presentation_title = self.build_pptx(
                rendered_markdown,
                file_name=filename,
                user_name=user_name,
            )
        except Exception as exc:
            self._debug_log(
                "PPTX generation failed",
                error=str(exc),
                message_id=message_id,
                filename=filename,
            )
            await self.emit_error(f"PPTX export failed: {exc}", __event_emitter__)
            await self.emit_status("PPTX generation failed.", True, __event_emitter__)
            return {
                "content": f"PPTX export failed: {exc}",
                "mermaid_extract_result": extract_result,
                "mermaid_extract_error": extract_error,
            }

        await self.emit_status("Starting download...", False, __event_emitter__)
        result = await self.download_file(
            pptx_bytes=pptx_bytes,
            filename=filename,
            __event_emitter__=__event_emitter__,
            __event_call__=__event_call__,
        )

        self._debug_log(
            "Action completed successfully",
            filename=filename,
            slide_count=slide_count,
            presentation_title=presentation_title,
            size_bytes=len(pptx_bytes),
            download_result=result,
        )
        await self.emit_status("PPTX export complete.", True, __event_emitter__)
        return {
            "content": f"Exported message to PPTX: {filename}",
            "result": result,
            "presentation_title": presentation_title,
            "slide_count": slide_count,
            "mermaid_diagrams_embedded": len(diagrams),
            "mermaid_extract_result": extract_result,
            "mermaid_extract_error": extract_error,
        }
