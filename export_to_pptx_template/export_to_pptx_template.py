"""
title: Export to PPTX (Template)
description: Export Assistant Message to PPTX using a template-aware LLM pipeline.
version: 1.0.0
author: Helmi Chaouachi
git_url: https://github.com/Helmi97/open-webui-extensions/tree/main/export_to_pptx_template
icon_url: https://www.svgrepo.com/show/373991/powerpoint2.svg
required_open_webui_version: 0.8.0
requirements: python-pptx
"""

from __future__ import annotations

import base64
import io
import json
import logging
import os
import re
import tempfile
import time
import urllib.parse
import urllib.request
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pydantic import BaseModel, Field
from pptx import Presentation

from open_webui.models.users import UserModel
from open_webui.utils.chat import generate_chat_completion

LOGGER = logging.getLogger(__name__)

DEFAULT_PREPROCESSING_PROMPT = """
You are preparing raw source text so it is easier to turn into a slide deck.

Task:
- Rewrite the input into a Marp-style markdown presentation.
- Use `---` between slides.
- Keep slides concise.
- Prefer short slide titles.
- Prefer bullet points over dense paragraphs.
- Preserve the meaning of the source.
- Split overloaded content into multiple slides when needed.
- Do not add speaker notes.
- Do not use code fences.
- Return plain markdown only.
""".strip()

DEFAULT_PROCESSING_PROMPT = """
You are converting presentation content into a structured slide plan for a
PowerPoint template.

You will receive:
1. The presentation text.
2. The available PowerPoint layouts, including raw placeholder metadata.

Your job:
- Choose only from the provided layout names.
- Build a slide plan in JSON.
- Keep slide text concise and presentation-friendly.
- You may split or merge content when needed to fit the template better.
- Do not invent layout names.
- Do not emit placeholder indices or coordinates in the output.
- Return valid JSON only.

Expected JSON shape:
{
  "slides": [
    {
      "layout_name": "Title Slide",
      "title": "string",
      "subtitle": "string",
      "body": "string",
      "bullets": ["string", "string"],
      "left_title": "string",
      "left_bullets": ["string"],
      "right_title": "string",
      "right_bullets": ["string"],
      "notes": "optional debugging or fallback note"
    }
  ]
}

Rules:
- Omit fields you do not need.
- Use "title" for the slide title whenever possible.
- Use "bullets" for normal content slides.
- Use left/right fields only when the content clearly has two parallel sections.
- Return JSON only, with no markdown fences and no extra commentary.
""".strip()

DEFAULT_POSTPROCESSING_PROMPT = """
You are validating and repairing a PowerPoint slide plan JSON document.

Task:
- Keep the same general meaning and slide order.
- Ensure all slides use only the provided layout names.
- Shorten titles and bullets when they are too long.
- Remove empty or redundant fields.
- Prefer concise presentation text.
- Return valid JSON only.

Do not add commentary.
Do not add markdown fences.
""".strip()


@dataclass(frozen=True, slots=True)
class LayoutMetadata:
    master_index: int
    layout_index: int
    layout_name: str
    placeholders: list[dict[str, Any]]


class _ManagedTemplateFile:
    def __init__(self, template: str | os.PathLike[str] | None, timeout_s: int):
        self.template = template
        self.timeout_s = timeout_s
        self._temp_path: str | None = None
        self.path: str | None = None

    def __enter__(self) -> str | None:
        if self.template is None:
            self.path = None
            return None

        template_str = str(self.template).strip()
        if not template_str:
            self.path = None
            return None

        if _is_url(template_str):
            self._temp_path = _download_template(template_str, self.timeout_s)
            self.path = self._temp_path
            return self.path

        local_path = Path(template_str).expanduser().resolve()
        if not local_path.exists():
            raise FileNotFoundError(f"Template not found: {local_path}")
        if local_path.suffix.lower() != ".pptx":
            raise ValueError(f"Template must be a .pptx file: {local_path}")

        self.path = str(local_path)
        return self.path

    def __exit__(self, exc_type, exc, tb) -> None:
        if self._temp_path and os.path.exists(self._temp_path):
            try:
                os.remove(self._temp_path)
            except OSError:
                pass


def _is_url(value: str) -> bool:
    parsed = urllib.parse.urlparse(value)
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def _download_template(url: str, timeout_s: int) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        temp_path = tmp.name

    try:
        with urllib.request.urlopen(url, timeout=timeout_s) as response:
            with open(temp_path, "wb") as output:
                output.write(response.read())
    except Exception:
        try:
            os.remove(temp_path)
        except OSError:
            pass
        raise

    return temp_path


class Action:
    class Valves(BaseModel):
        debug: bool = Field(
            default=False,
            description="Enable verbose server-side debug logging for this action.",
        )
        priority: int = Field(
            default=0,
            description="Controls button display order (lower = appears first).",
        )
        filename_prefix: str = Field(
            default="message",
            description="Prefix used for the downloaded PPTX file name.",
        )
        template: str = Field(
            default="",
            description=(
                "Optional PPTX template path or URL. Leave empty to use the default "
                "python-pptx presentation template."
            ),
        )
        template_request_timeout_s: int = Field(
            default=30,
            ge=1,
            le=300,
            description="HTTP timeout in seconds when downloading a remote template.",
        )
        llm_model: str = Field(
            default="",
            description=(
                "Optional model ID used for slide planning. Defaults to the current "
                "chat model."
            ),
        )
        llm_temperature: float = Field(
            default=0.2,
            ge=0.0,
            le=2.0,
            description="Temperature used for all slide-planning LLM calls.",
        )
        llm_max_tokens: int = Field(
            default=4096,
            ge=256,
            le=16384,
            description="Maximum tokens used for each slide-planning LLM call.",
        )
        enable_preprocessing: bool = Field(
            default=True,
            description=(
                "Run a first pass that rewrites the message into cleaner "
                "presentation-style markdown before slide planning."
            ),
        )
        enable_postprocessing: bool = Field(
            default=True,
            description=(
                "Run a final validation/repair pass on the generated slide plan JSON."
            ),
        )
        preprocessing_prompt: str = Field(
            default=DEFAULT_PREPROCESSING_PROMPT,
            description="System prompt used for the optional preprocessing pass.",
        )
        processing_prompt: str = Field(
            default=DEFAULT_PROCESSING_PROMPT,
            description="System prompt used for the required slide-planning pass.",
        )
        postprocessing_prompt: str = Field(
            default=DEFAULT_POSTPROCESSING_PROMPT,
            description="System prompt used for the optional repair pass.",
        )

    def __init__(self) -> None:
        self.valves = self.Valves()

    def _preview(self, value: Any, limit: int = 220) -> str:
        text = value if isinstance(value, str) else str(value)
        text = text.replace("\r", "\\r").replace("\n", "\\n")
        if len(text) > limit:
            return f"{text[:limit]}...<truncated {len(text) - limit} chars>"
        return text

    def _debug_log(self, message: str, **context: Any) -> None:
        if not self.valves.debug:
            return

        if not context:
            LOGGER.info("[export_to_pptx_template] %s", message)
            return

        rendered_parts: list[str] = []
        for key, value in context.items():
            if isinstance(value, (dict, list, tuple)):
                try:
                    rendered = json.dumps(value, default=str)
                except TypeError:
                    rendered = str(value)
            else:
                rendered = str(value)
            rendered_parts.append(f"{key}={self._preview(rendered)}")

        LOGGER.info(
            "[export_to_pptx_template] %s | %s",
            message,
            " | ".join(rendered_parts),
        )

    async def emit_status(
        self,
        description: str,
        done: bool = False,
        __event_emitter__=None,
    ) -> None:
        if __event_emitter__ is None:
            return

        await __event_emitter__(
            {
                "type": "status",
                "data": {"description": description, "done": done},
            }
        )

    async def emit_error(self, message: str, __event_emitter__=None) -> None:
        if __event_emitter__ is None:
            return

        await __event_emitter__(
            {
                "type": "notification",
                "data": {"type": "error", "content": message},
            }
        )

    def build_filename(self, message_id: str) -> str:
        safe_id = re.sub(r"[^A-Za-z0-9._-]+", "-", (message_id or "").strip())
        safe_id = safe_id.strip("-._") or "message"
        prefix = re.sub(r"[^A-Za-z0-9._-]+", "-", self.valves.filename_prefix).strip(
            "-._"
        )
        prefix = prefix or "message"
        return f"{prefix}-{safe_id}.pptx"

    def _extract_text(self, value: Any) -> str:
        if value is None:
            return ""

        if isinstance(value, str):
            return value

        if isinstance(value, list):
            parts = [self._extract_text(item).strip() for item in value]
            return "\n\n".join(part for part in parts if part)

        if isinstance(value, dict):
            item_type = value.get("type")
            if item_type in {"text", "input_text", "output_text"}:
                return self._extract_text(value.get("text"))

            if item_type == "message":
                return self._extract_text(value.get("content"))

            if "choices" in value:
                return self._extract_text(value.get("choices"))

            if "output" in value:
                return self._extract_text(value.get("output"))

            if "message" in value:
                return self._extract_text(value.get("message"))

            for key in ("content", "output_text", "text", "value"):
                if key in value:
                    return self._extract_text(value.get(key))

            return ""

        return str(value)

    def _extract_chat_completion_text(self, response: Any) -> str:
        if response is None:
            return ""

        if isinstance(response, dict):
            choices = response.get("choices")
            if isinstance(choices, list) and choices:
                first_choice = choices[0]
                if isinstance(first_choice, dict):
                    message = first_choice.get("message")
                    if message is not None:
                        return self._extract_text(message).strip()

            output = response.get("output")
            if output is not None:
                extracted_output = self._extract_text(output).strip()
                if extracted_output:
                    return extracted_output

            for key in ("output_text", "content", "text"):
                if key in response:
                    extracted = self._extract_text(response.get(key)).strip()
                    if extracted:
                        return extracted

        response_body = getattr(response, "body", None)
        if response_body:
            try:
                payload = json.loads(response_body.decode("utf-8"))
            except Exception:
                return ""
            return self._extract_chat_completion_text(payload)

        return ""

    def _log_selected_message(
        self,
        source: str,
        message_text: str,
        message_id: str | None = None,
    ) -> None:
        log_context: dict[str, Any] = {
            "length": len(message_text),
            "preview": self._preview(message_text),
        }
        if message_id:
            log_context["message_id"] = message_id
        self._debug_log(source, **log_context)

    def _extract_message_text(self, body: dict) -> str:
        direct_text = self._extract_text(body.get("content")).strip()
        if direct_text:
            self._log_selected_message(
                "Selected message from body.content",
                direct_text,
                body.get("id"),
            )
            return direct_text

        message = body.get("message")
        if isinstance(message, dict):
            message_text = self._extract_text(message.get("content")).strip()
            if message_text:
                self._log_selected_message(
                    "Selected message from body.message.content",
                    message_text,
                    message.get("id"),
                )
                return message_text

        target_id = body.get("id")
        messages = body.get("messages")
        if isinstance(messages, list):
            if target_id:
                for message in messages:
                    if not isinstance(message, dict):
                        continue

                    candidate_ids = [
                        message.get("id"),
                        message.get("message_id"),
                        (
                            (message.get("info") or {}).get("id")
                            if isinstance(message.get("info"), dict)
                            else None
                        ),
                        (
                            (message.get("meta") or {}).get("id")
                            if isinstance(message.get("meta"), dict)
                            else None
                        ),
                    ]

                    if target_id not in candidate_ids:
                        continue

                    message_text = self._extract_text(message.get("content")).strip()
                    if message_text:
                        self._log_selected_message(
                            "Selected message from body.messages by id",
                            message_text,
                            target_id,
                        )
                        return message_text

            for message in reversed(messages):
                if not isinstance(message, dict) or message.get("role") != "assistant":
                    continue

                message_text = self._extract_text(message.get("content")).strip()
                if message_text:
                    self._log_selected_message(
                        "Selected latest assistant message from body.messages",
                        message_text,
                        message.get("id"),
                    )
                    return message_text

        history = body.get("history")
        if isinstance(history, dict):
            history_messages = history.get("messages", {})
            current_id = target_id or history.get("currentId")

            if current_id and current_id in history_messages:
                current_message = history_messages[current_id]
                if isinstance(current_message, dict):
                    message_text = self._extract_text(
                        current_message.get("content")
                    ).strip()
                    if message_text:
                        self._log_selected_message(
                            "Selected message from body.history current message",
                            message_text,
                            str(current_id),
                        )
                        return message_text

            if isinstance(history_messages, dict):
                for message in reversed(list(history_messages.values())):
                    if not isinstance(message, dict) or message.get("role") != "assistant":
                        continue

                    message_text = self._extract_text(message.get("content")).strip()
                    if message_text:
                        self._log_selected_message(
                            "Selected latest assistant message from body.history",
                            message_text,
                            message.get("id"),
                        )
                        return message_text

        self._debug_log(
            "No assistant message content found during extraction",
            body_keys=sorted(body.keys()),
        )
        return ""

    def _get_user_model(self, user_data: Any) -> UserModel | None:
        if isinstance(user_data, UserModel):
            return user_data

        if not isinstance(user_data, dict) or not user_data:
            return None

        try:
            return UserModel(**user_data)
        except Exception:
            timestamp = int(time.time())
            fallback_user_data = {
                "id": user_data.get("id", ""),
                "email": user_data.get("email", ""),
                "username": user_data.get("username"),
                "role": user_data.get("role", "user"),
                "name": user_data.get("name") or user_data.get("email") or "User",
                "profile_image_url": user_data.get("profile_image_url"),
                "profile_banner_image_url": user_data.get(
                    "profile_banner_image_url"
                ),
                "bio": user_data.get("bio"),
                "gender": user_data.get("gender"),
                "date_of_birth": user_data.get("date_of_birth"),
                "timezone": user_data.get("timezone"),
                "presence_state": user_data.get("presence_state"),
                "status_emoji": user_data.get("status_emoji"),
                "status_message": user_data.get("status_message"),
                "status_expires_at": user_data.get("status_expires_at"),
                "info": user_data.get("info"),
                "settings": user_data.get("settings"),
                "oauth": user_data.get("oauth"),
                "scim": user_data.get("scim"),
                "last_active_at": user_data.get("last_active_at", timestamp),
                "updated_at": user_data.get("updated_at", timestamp),
                "created_at": user_data.get("created_at", timestamp),
            }
            try:
                return UserModel(**fallback_user_data)
            except Exception:
                return None

    def _resolve_current_model_id(self, body: dict, __model__) -> str:
        return (
            self.valves.llm_model.strip()
            or body.get("model")
            or (__model__ or {}).get("id")
            or ""
        )

    async def _call_chat_completion(
        self,
        *,
        model_id: str,
        messages: list[dict[str, str]],
        temperature: float,
        max_tokens: int,
        __request__,
        user_model: UserModel,
        response_format: dict[str, Any] | None = None,
    ) -> str:
        payload: dict[str, Any] = {
            "model": model_id,
            "stream": False,
            "temperature": temperature,
            "max_tokens": max_tokens,
            "messages": messages,
        }
        if response_format is not None:
            payload["response_format"] = response_format

        self._debug_log(
            "Requesting Open WebUI chat completion",
            model_id=model_id,
            temperature=temperature,
            max_tokens=max_tokens,
            message_count=len(messages),
            response_format=response_format,
        )

        try:
            response = await generate_chat_completion(
                __request__,
                payload,
                user_model,
                bypass_system_prompt=True,
            )
            content = self._extract_chat_completion_text(response).strip()
            if content:
                return content
        except Exception as exc:
            if response_format is None:
                raise

            self._debug_log(
                "Chat completion with response_format failed; retrying without it",
                error=str(exc),
                model_id=model_id,
            )

        if response_format is not None:
            retry_payload = dict(payload)
            retry_payload.pop("response_format", None)
            response = await generate_chat_completion(
                __request__,
                retry_payload,
                user_model,
                bypass_system_prompt=True,
            )
            content = self._extract_chat_completion_text(response).strip()
            if content:
                return content

        raise ValueError("Chat completion returned empty content.")

    async def _run_preprocessing(
        self,
        *,
        text: str,
        layouts: list[LayoutMetadata],
        model_id: str,
        user_model: UserModel,
        __request__,
    ) -> str:
        messages = [
            {"role": "system", "content": self.valves.preprocessing_prompt.strip()},
            {
                "role": "user",
                "content": json.dumps(
                    {
                        "text": text,
                        "available_layouts": self._layouts_to_jsonable(layouts),
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
            },
        ]
        content = await self._call_chat_completion(
            model_id=model_id,
            messages=messages,
            temperature=self.valves.llm_temperature,
            max_tokens=self.valves.llm_max_tokens,
            __request__=__request__,
            user_model=user_model,
        )
        if not content:
            raise ValueError("Preprocessing returned empty content.")
        return content

    async def _run_processing(
        self,
        *,
        text: str,
        layouts: list[LayoutMetadata],
        model_id: str,
        user_model: UserModel,
        __request__,
    ) -> dict[str, Any]:
        messages = [
            {"role": "system", "content": self.valves.processing_prompt.strip()},
            {
                "role": "user",
                "content": json.dumps(
                    {
                        "text": text,
                        "available_layouts": self._layouts_to_jsonable(layouts),
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
            },
        ]
        content = await self._call_chat_completion(
            model_id=model_id,
            messages=messages,
            temperature=self.valves.llm_temperature,
            max_tokens=self.valves.llm_max_tokens,
            __request__=__request__,
            user_model=user_model,
            response_format={"type": "json_object"},
        )
        slide_plan = self._parse_json_object(content)
        self._validate_slide_plan_shape(slide_plan)
        return slide_plan

    async def _run_postprocessing(
        self,
        *,
        slide_plan: dict[str, Any],
        layouts: list[LayoutMetadata],
        model_id: str,
        user_model: UserModel,
        __request__,
    ) -> dict[str, Any]:
        messages = [
            {"role": "system", "content": self.valves.postprocessing_prompt.strip()},
            {
                "role": "user",
                "content": json.dumps(
                    {
                        "slide_plan": slide_plan,
                        "available_layouts": self._layouts_to_jsonable(layouts),
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
            },
        ]
        content = await self._call_chat_completion(
            model_id=model_id,
            messages=messages,
            temperature=self.valves.llm_temperature,
            max_tokens=self.valves.llm_max_tokens,
            __request__=__request__,
            user_model=user_model,
            response_format={"type": "json_object"},
        )
        repaired_plan = self._parse_json_object(content)
        self._validate_slide_plan_shape(repaired_plan)
        return repaired_plan

    async def _text_to_pptx(
        self,
        *,
        text: str,
        template: str | os.PathLike[str] | None,
        model_id: str,
        user_model: UserModel,
        __request__,
    ) -> tuple[Presentation, dict[str, Any]]:
        if not isinstance(text, str) or not text.strip():
            raise ValueError("`text` must be a non-empty string.")

        with _ManagedTemplateFile(
            template,
            timeout_s=self.valves.template_request_timeout_s,
        ) as template_path:
            presentation = self._load_presentation(template_path)
            layouts = self._extract_layout_metadata(presentation)

            self._debug_log(
                "Loaded presentation template",
                template_path=template_path,
                layout_count=len(layouts),
            )

            working_text = text
            if self.valves.enable_preprocessing:
                working_text = await self._run_preprocessing(
                    text=text,
                    layouts=layouts,
                    model_id=model_id,
                    user_model=user_model,
                    __request__=__request__,
                )

            slide_plan = await self._run_processing(
                text=working_text,
                layouts=layouts,
                model_id=model_id,
                user_model=user_model,
                __request__=__request__,
            )

            if self.valves.enable_postprocessing:
                slide_plan = await self._run_postprocessing(
                    slide_plan=slide_plan,
                    layouts=layouts,
                    model_id=model_id,
                    user_model=user_model,
                    __request__=__request__,
                )

            self._render_slide_plan(presentation, slide_plan, layouts)
            return presentation, slide_plan

    def _load_presentation(self, template_path: str | None) -> Presentation:
        if template_path is None:
            return Presentation()
        return Presentation(template_path)

    def _extract_layout_metadata(
        self,
        presentation: Presentation,
    ) -> list[LayoutMetadata]:
        layouts: list[LayoutMetadata] = []

        for master_index, master in enumerate(presentation.slide_masters):
            for layout_index, layout in enumerate(master.slide_layouts):
                placeholder_rows: list[dict[str, Any]] = []

                for placeholder in layout.placeholders:
                    pformat = placeholder.placeholder_format
                    placeholder_rows.append(
                        {
                            "idx": pformat.idx,
                            "type": self._enum_name(pformat.type),
                            "name": placeholder.name,
                            "left": int(placeholder.left),
                            "top": int(placeholder.top),
                            "width": int(placeholder.width),
                            "height": int(placeholder.height),
                        }
                    )

                layouts.append(
                    LayoutMetadata(
                        master_index=master_index,
                        layout_index=layout_index,
                        layout_name=layout.name,
                        placeholders=sorted(
                            placeholder_rows,
                            key=lambda row: row["idx"],
                        ),
                    )
                )

        return layouts

    def _enum_name(self, value: Any) -> str:
        return getattr(value, "name", str(value))

    def _layouts_to_jsonable(
        self,
        layouts: Iterable[LayoutMetadata],
    ) -> list[dict[str, Any]]:
        return [
            {
                "master_index": layout.master_index,
                "layout_index": layout.layout_index,
                "layout_name": layout.layout_name,
                "placeholders": layout.placeholders,
            }
            for layout in layouts
        ]

    def _parse_json_object(self, content: str) -> dict[str, Any]:
        stripped = content.strip()
        fence_match = re.fullmatch(r"```(?:json)?\s*(.*?)\s*```", stripped, re.DOTALL)
        if fence_match:
            stripped = fence_match.group(1).strip()

        try:
            obj = json.loads(stripped)
        except json.JSONDecodeError as exc:
            raise ValueError("Model did not return valid JSON.") from exc

        if not isinstance(obj, dict):
            raise ValueError("Top-level JSON must be an object.")
        return obj

    def _validate_slide_plan_shape(self, slide_plan: dict[str, Any]) -> None:
        slides = slide_plan.get("slides")
        if not isinstance(slides, list):
            raise ValueError("Slide plan JSON must contain a 'slides' list.")

        for index, slide in enumerate(slides):
            if not isinstance(slide, dict):
                raise ValueError(f"Slide {index} must be a JSON object.")
            layout_name = slide.get("layout_name")
            if not isinstance(layout_name, str) or not layout_name.strip():
                raise ValueError(
                    f"Slide {index} must include a non-empty 'layout_name'."
                )

    def _resolve_layout_identity(
        self,
        requested_layout_name: str,
        layout_index_by_name: dict[str, tuple[int, int]],
    ) -> tuple[int, int]:
        if requested_layout_name in layout_index_by_name:
            return layout_index_by_name[requested_layout_name]

        normalized = requested_layout_name.casefold()
        for layout_name, identity in layout_index_by_name.items():
            if layout_name.casefold() == normalized:
                return identity

        first_available = next(iter(layout_index_by_name.values()), None)
        if first_available is None:
            raise ValueError("No slide layouts are available in the presentation.")

        self._debug_log(
            "Unknown layout in slide plan; falling back to first layout",
            requested_layout_name=requested_layout_name,
        )
        return first_available

    def _render_slide_plan(
        self,
        presentation: Presentation,
        slide_plan: dict[str, Any],
        layouts: list[LayoutMetadata],
    ) -> None:
        layout_index_by_name = {
            layout.layout_name: (layout.master_index, layout.layout_index)
            for layout in layouts
        }

        slides = slide_plan.get("slides", [])
        for slide_data in slides:
            requested_layout = str(slide_data.get("layout_name", "")).strip()
            master_index, layout_index = self._resolve_layout_identity(
                requested_layout,
                layout_index_by_name,
            )
            slide_layout = presentation.slide_masters[master_index].slide_layouts[
                layout_index
            ]
            slide = presentation.slides.add_slide(slide_layout)
            self._populate_slide(slide, slide_data)

    def _populate_slide(self, slide: Any, slide_data: dict[str, Any]) -> None:
        title_text = self._as_clean_text(slide_data.get("title"))
        subtitle_text = self._as_clean_text(slide_data.get("subtitle"))
        body_text = self._as_clean_text(slide_data.get("body"))
        bullets = self._as_clean_list(slide_data.get("bullets"))
        left_title = self._as_clean_text(slide_data.get("left_title"))
        left_bullets = self._as_clean_list(slide_data.get("left_bullets"))
        right_title = self._as_clean_text(slide_data.get("right_title"))
        right_bullets = self._as_clean_list(slide_data.get("right_bullets"))

        if title_text and getattr(slide.shapes, "title", None) is not None:
            try:
                slide.shapes.title.text = title_text
            except Exception:
                pass

        non_title_placeholders = self._get_non_title_placeholders(slide)

        if (
            left_title
            or left_bullets
            or right_title
            or right_bullets
        ) and len(non_title_placeholders) >= 2:
            left_text = self._combine_heading_and_bullets(left_title, left_bullets)
            right_text = self._combine_heading_and_bullets(right_title, right_bullets)
            if left_text:
                self._write_text_to_placeholder(non_title_placeholders[0], left_text)
            if right_text:
                self._write_text_to_placeholder(non_title_placeholders[1], right_text)
            return

        if subtitle_text and non_title_placeholders:
            self._write_text_to_placeholder(non_title_placeholders[0], subtitle_text)
            if len(non_title_placeholders) >= 2:
                main_text = self._combine_body_and_bullets(body_text, bullets)
                if main_text:
                    self._write_text_to_placeholder(non_title_placeholders[1], main_text)
                return

        if non_title_placeholders:
            main_text = self._combine_body_and_bullets(body_text, bullets)
            if main_text:
                self._write_text_to_placeholder(non_title_placeholders[0], main_text)

    def _get_non_title_placeholders(self, slide: Any) -> list[Any]:
        placeholders: list[Any] = []
        for shape in slide.placeholders:
            if not getattr(shape, "is_placeholder", False):
                continue

            placeholder_type = self._enum_name(shape.placeholder_format.type)
            if placeholder_type in {"TITLE", "CENTER_TITLE"}:
                continue

            placeholders.append(shape)

        return sorted(placeholders, key=lambda shape: (int(shape.top), int(shape.left)))

    def _write_text_to_placeholder(self, placeholder: Any, text: str) -> None:
        if not text or not getattr(placeholder, "has_text_frame", False):
            return

        placeholder.text = text

    def _as_clean_text(self, value: Any) -> str:
        if value is None:
            return ""
        return str(value).strip()

    def _as_clean_list(self, value: Any) -> list[str]:
        if value is None:
            return []
        if isinstance(value, str):
            value = [value]
        if not isinstance(value, list):
            return []

        cleaned: list[str] = []
        for item in value:
            item_text = self._as_clean_text(item)
            if item_text:
                cleaned.append(item_text)
        return cleaned

    def _combine_body_and_bullets(self, body: str, bullets: list[str]) -> str:
        parts: list[str] = []
        if body:
            parts.append(body)
        if bullets:
            parts.append("\n".join(f"- {bullet}" for bullet in bullets))
        return "\n\n".join(parts).strip()

    def _combine_heading_and_bullets(self, heading: str, bullets: list[str]) -> str:
        parts: list[str] = []
        if heading:
            parts.append(heading)
        if bullets:
            parts.append("\n".join(f"- {bullet}" for bullet in bullets))
        return "\n\n".join(parts).strip()

    def _guess_presentation_title(
        self,
        slide_plan: dict[str, Any],
        message_id: str,
    ) -> str:
        slides = slide_plan.get("slides")
        if isinstance(slides, list):
            for slide in slides:
                if not isinstance(slide, dict):
                    continue
                title = self._as_clean_text(slide.get("title"))
                if title:
                    return title
        return message_id

    async def download_file(
        self,
        pptx_bytes: bytes,
        filename: str,
        __event_emitter__=None,
        __event_call__=None,
    ) -> Any:
        encoded = base64.b64encode(pptx_bytes).decode("ascii")
        js_code = f"""
const base64 = {json.dumps(encoded)};
const filename = {json.dumps(filename)};
const binary = atob(base64);
const bytes = new Uint8Array(binary.length);

for (let i = 0; i < binary.length; i++) {{
  bytes[i] = binary.charCodeAt(i);
}}

const blob = new Blob(
  [bytes],
  {{ type: "application/vnd.openxmlformats-officedocument.presentationml.presentation" }}
);
const url = URL.createObjectURL(blob);

try {{
  const a = document.createElement("a");
  a.href = url;
  a.download = filename;
  a.style.display = "none";
  document.body.appendChild(a);
  a.click();
  a.remove();
}} finally {{
  setTimeout(() => URL.revokeObjectURL(url), 4000);
}}

return {{ success: true, filename, size: bytes.length }};
"""

        payload = {"type": "execute", "data": {"code": js_code}}

        if __event_call__ is not None:
            return await __event_call__(payload)

        if __event_emitter__ is not None:
            await __event_emitter__(payload)
            return {"success": True, "filename": filename, "size": len(pptx_bytes)}

        return None

    async def action(
        self,
        body: dict,
        __user__=None,
        __event_emitter__=None,
        __event_call__=None,
        __request__=None,
        __model__=None,
        **kwargs,
    ) -> dict[str, Any]:
        self._debug_log(
            "Action invoked",
            body_keys=sorted(body.keys()),
            chat_id=body.get("chat_id"),
            message_id=body.get("id"),
            current_model=body.get("model"),
            model_override=self.valves.llm_model,
            template=self.valves.template,
        )

        message_id = str(body.get("id") or "").strip()
        if not message_id:
            await self.emit_error(
                "PPTX export failed: could not determine the current message id.",
                __event_emitter__,
            )
            return {
                "content": "Could not determine the current message id from body['id']."
            }

        filename = self.build_filename(message_id)

        await self.emit_status("Reading message content...", False, __event_emitter__)
        message_text = self._extract_message_text(body)
        if not message_text.strip():
            await self.emit_error(
                "PPTX export failed: no assistant message content found.",
                __event_emitter__,
            )
            await self.emit_status(
                "No assistant message content found.",
                True,
                __event_emitter__,
            )
            return {"content": "No assistant message content found."}

        if __request__ is None:
            await self.emit_error(
                "PPTX export failed: __request__ is required for Open WebUI chat completion.",
                __event_emitter__,
            )
            await self.emit_status("PPTX export failed.", True, __event_emitter__)
            return {
                "content": "PPTX export failed: __request__ is unavailable for chat completion."
            }

        user_model = self._get_user_model(__user__)
        if user_model is None:
            await self.emit_error(
                "PPTX export failed: could not build a valid Open WebUI user context.",
                __event_emitter__,
            )
            await self.emit_status("PPTX export failed.", True, __event_emitter__)
            return {
                "content": "PPTX export failed: could not build a valid Open WebUI user context."
            }

        model_id = self._resolve_current_model_id(body, __model__)
        if not model_id:
            await self.emit_error(
                "PPTX export failed: no model id is available for slide planning.",
                __event_emitter__,
            )
            await self.emit_status("PPTX export failed.", True, __event_emitter__)
            return {"content": "PPTX export failed: no model id is available."}

        template = (self.valves.template or "").strip() or None

        try:
            await self.emit_status(
                "Planning slides with the current chat model...",
                False,
                __event_emitter__,
            )
            presentation, slide_plan = await self._text_to_pptx(
                text=message_text,
                template=template,
                model_id=model_id,
                user_model=user_model,
                __request__=__request__,
            )
        except Exception as exc:
            self._debug_log(
                "PPTX planning/generation failed",
                error=str(exc),
                model_id=model_id,
                template=template,
            )
            await self.emit_error(f"PPTX export failed: {exc}", __event_emitter__)
            await self.emit_status("PPTX export failed.", True, __event_emitter__)
            return {"content": f"PPTX export failed: {exc}"}

        await self.emit_status("Building PPTX...", False, __event_emitter__)

        output = io.BytesIO()
        presentation.save(output)
        pptx_bytes = output.getvalue()
        slide_count = len(presentation.slides)
        presentation_title = self._guess_presentation_title(slide_plan, message_id)

        self._debug_log(
            "Built PPTX bytes",
            filename=filename,
            size_bytes=len(pptx_bytes),
            slide_count=slide_count,
            presentation_title=presentation_title,
        )

        await self.emit_status("Starting download...", False, __event_emitter__)
        result = await self.download_file(
            pptx_bytes=pptx_bytes,
            filename=filename,
            __event_emitter__=__event_emitter__,
            __event_call__=__event_call__,
        )

        await self.emit_status("PPTX export complete.", True, __event_emitter__)
        return {
            "content": f"Exported message to PPTX: {filename}",
            "result": result,
            "presentation_title": presentation_title,
            "slide_count": slide_count,
            "template": template,
        }
